import streamlit as st
import pandas as pd
from datetime import datetime, date, timedelta
import json
import plotly.graph_objects as go
import plotly.express as px
# import numpy as np
import xml.etree.ElementTree as ET
import io
from io import StringIO
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import tempfile
# import base64
import matplotlib.pyplot as plt
import tabulate
from PIL import Image

# Import Snowflake-related libraries

try:
    from snowflake.snowpark.context import get_active_session
    from snowflake.snowpark.functions import col
    SNOWFLAKE_ENABLED = True
except ImportError:
    SNOWFLAKE_ENABLED = False
    st.sidebar.error("Snowflake libraries not found. Database features will be disabled.")
# Set page configuration
st.set_page_config(
    page_title="Snowflake Customer Sizing Tool",
    page_icon="❄️",
    layout="wide"
)

# Custom CSS for better styling

st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        color: #29B5E8;
        text-align: center;
        margin-bottom: 1rem;
    }
    .section-header {
        font-size: 1.5rem;
        color: #29B5E8;
        margin-top: 1rem;
        margin-bottom: 0.5rem;
        background-color: #f0f8ff;
        padding: 0.5rem;
        border-radius: 5px;
    }
    .category-label {
        font-weight: bold;
        color: #0066cc;
    }
    .stButton>button {
        background-color: #29B5E8;
        color: white;
    }
    .summary-box {
        background-color: #f0f8ff;
        padding: 1rem;
        border-radius: 5px;
        margin-bottom: 1rem;
    }
    .svg-container {
        border: 2px solid #29B5E8;
        border-radius: 10px;
        padding: 10px 0;
        background-color: #f8f9fa;
    }

    /* ONLY TARGET TABLE INPUTS - NOT ALL INPUTS */

    /* Table text inputs only */
    div[data-testid*="wh_warehouse_name"] input {
        font-size: 11px !important;
        font-family: "Source Sans Pro", sans-serif !important;
        padding: 3px 6px !important;
        height: 32px !important;
        line-height: 1.2 !important;
        width: 200px !important;
        min-width: 200px !important;
    }

    /* Table number inputs only */
    div[data-testid*="wh_days_per_month"] input,
    div[data-testid*="wh_hours_per_day"] input,
    div[data-testid*="wh_warehouses_when_active"] input,
    div[data-testid*="wh_annual_growth_rate"] input,
    div[data-testid*="wh_dev_start_month"] input,
    div[data-testid*="wh_go_live_month"] input,
    div[data-testid*="wh_monthly_credits"] input {
        font-size: 11px !important;
        font-family: "Source Sans Pro", sans-serif !important;
        padding: 3px 4px !important;
        height: 32px !important;
        line-height: 1.2 !important;
        width: 50px !important;
        max-width: 50px !important;
        min-width: 50px !important;
        text-align: center !important;
    }

    /* Table selectboxes only - target by specific test IDs */
    div[data-testid*="wh_category"] div[data-baseweb="select"],
    div[data-testid*="wh_feature"] div[data-baseweb="select"],
    div[data-testid*="wh_size"] div[data-baseweb="select"],
    div[data-testid*="wh_ramp_curve"] div[data-baseweb="select"] {
        min-height: 32px !important;
        height: auto !important;
    }

    /* Table dropdown display text - small font */
    div[data-testid*="wh_category"] div[data-baseweb="select"] div,
    div[data-testid*="wh_category"] div[data-baseweb="select"] span,
    div[data-testid*="wh_feature"] div[data-baseweb="select"] div,
    div[data-testid*="wh_feature"] div[data-baseweb="select"] span,
    div[data-testid*="wh_size"] div[data-baseweb="select"] div,
    div[data-testid*="wh_size"] div[data-baseweb="select"] span,
    div[data-testid*="wh_ramp_curve"] div[data-baseweb="select"] div,
    div[data-testid*="wh_ramp_curve"] div[data-baseweb="select"] span {
        font-size: 11px !important;
        font-family: "Source Sans Pro", sans-serif !important;
    }

    /* Table dropdown menu items when opened */
    div[data-testid*="wh_category"] div[role="listbox"] div,
    div[data-testid*="wh_category"] div[role="option"],
    div[data-testid*="wh_feature"] div[role="listbox"] div,
    div[data-testid*="wh_feature"] div[role="option"],
    div[data-testid*="wh_size"] div[role="listbox"] div,
    div[data-testid*="wh_size"] div[role="option"],
    div[data-testid*="wh_ramp_curve"] div[role="listbox"] div,
    div[data-testid*="wh_ramp_curve"] div[role="option"] {
        font-size: 11px !important;
        font-family: "Source Sans Pro", sans-serif !important;
    }

    /* Table dropdown widths */
    div[data-testid*="wh_category"] div[data-baseweb="select"] {
        min-width: 120px !important;
        width: 120px !important;
    }

    div[data-testid*="wh_feature"] div[data-baseweb="select"] {
        min-width: 160px !important;
        width: 160px !important;
    }

    div[data-testid*="wh_size"] div[data-baseweb="select"] {
        min-width: 70px !important;
        width: 70px !important;
    }

    div[data-testid*="wh_ramp_curve"] div[data-baseweb="select"] {
        min-width: 110px !important;
        width: 110px !important;
    }

    /* TABLE HEADERS */
    .warehouse-table-header {
        font-size: 10px !important;
        font-weight: bold !important;
        padding: 2px !important;
        margin: 0 !important;
        text-align: center !important;
        line-height: 1.1 !important;
    }

    /* TABLE BUTTONS only */
    div[data-testid*="wh_delete"] button {
        height: 32px !important;
        padding: 2px 8px !important;
        font-size: 11px !important;
        min-width: 32px !important;
    }

    /* REDUCE COLUMN PADDING for table only */
    .stColumn > div {
        padding-left: 0.1rem !important;
        padding-right: 0.1rem !important;
    }

    .element-container {
        margin-bottom: 0.1rem !important;
    }
</style>
""", unsafe_allow_html=True)

# App title and introduction

def read_svg(path):
    # Check if file exists before trying to read
    if not os.path.exists(path):
        return None
    with open(path, 'r') as f:
        svg_string = f.read()
    return svg_string

# Assuming Snowflake_Logo.svg is in the same directory.
# If not, this will be handled gracefully.

svg_content = read_svg("Snowflake_Logo.svg")
if svg_content:
    st.sidebar.image(svg_content, width=200)

# Helper Functions

def get_recommended_warehouse_config():
    """
    Generate recommended warehouse configurations based on collected form data
    """
    recommended_configs = []

    # Get form data for recommendations
    pipeline_count = int(st.session_state.form_data.get('pipeline_count', 0))
    analytics_count = int(st.session_state.form_data.get('analytics_workload_count', 0))
    other_count = int(st.session_state.form_data.get('other_workload_count', 0)) if st.session_state.form_data.get('has_other_workloads', 'No') == 'Yes' else 0

    # Add ETL/ELT warehouses based on pipeline data
    for i in range(pipeline_count):
        pipeline_name = st.session_state.form_data.get(f'pipeline_{i}_name', f'Pipeline {i+1}')
        concurrent_jobs = int(st.session_state.form_data.get(f'pipeline_{i}_concurrent_jobs', 5))
        volume_per_job = float(st.session_state.form_data.get(f'pipeline_{i}_volume_per_job', 10.0))

        # Determine warehouse size based on workload
        workload_intensity = volume_per_job * concurrent_jobs
        if workload_intensity < 50:
            size = "XS"
        elif workload_intensity < 200:
            size = "S"
        elif workload_intensity < 500:
            size = "M"
        elif workload_intensity < 1000:
            size = "L"
        else:
            size = "XL"

        recommended_configs.append({
            'category': 'Warehouse',
            'feature': '',
            'monthly_credits': 0,
            'warehouse_name': f'ETL_{pipeline_name.replace(" ", "_")}',
            'size': size,
            'days_per_month': int(st.session_state.form_data.get(f'pipeline_{i}_days_per_month', 22)),
            'hours_per_day': 8,  # Reasonable default for ETL
            'warehouses_when_active': 1,
            'annual_growth_rate': 20,
            'dev_start_month': 1,
            'go_live_month': 6,
            'ramp_curve': 'Linear',
            'total': 0
        })

    # Add Analytics warehouses based on analytics workload data
    for i in range(analytics_count):
        analytics_name = st.session_state.form_data.get(f'analytics_{i}_name', f'Analytics {i+1}')
        concurrent_queries = int(st.session_state.form_data.get(f'analytics_concurrent_queries', 20))
        power_users_pct = int(st.session_state.form_data.get(f'analytics_{i}_power_users', 10))

        # Determine warehouse size based on complexity
        complexity_factor = concurrent_queries * (1 + power_users_pct/100)
        if complexity_factor < 15:
            size = "XS"
        elif complexity_factor < 40:
            size = "S"
        elif complexity_factor < 80:
            size = "M"
        elif complexity_factor < 150:
            size = "L"
        else:
            size = "XL"

        recommended_configs.append({
            'category': 'Warehouse',
            'feature': '',
            'monthly_credits': 0,
            'warehouse_name': f'ANALYTICS_{analytics_name.replace(" ", "_")}',
            'size': size,
            'days_per_month': int(st.session_state.form_data.get(f'analytics_{i}_days_per_month', 22)),
            'hours_per_day': int(st.session_state.form_data.get(f'analytics_{i}_hours_per_day', 8)),
            'warehouses_when_active': 1,
            'annual_growth_rate': 20,
            'dev_start_month': 1,
            'go_live_month': 6,
            'ramp_curve': 'Linear',
            'total': 0
        })

    # Add Other workload warehouses
    for i in range(other_count):
        other_name = st.session_state.form_data.get(f'other_{i}_name', f'Other {i+1}')
        other_type = st.session_state.form_data.get(f'other_{i}_type', 'Other')
        concurrent_queries = int(st.session_state.form_data.get(f'other_{i}_concurrent_queries', 10))
        data_volume = float(st.session_state.form_data.get(f'other_{i}_data_volume', 5.0))

        # Determine warehouse size
        workload_factor = concurrent_queries * (data_volume / 10)
        if workload_factor < 10:
            size = "XS"
        elif workload_factor < 30:
            size = "S"
        elif workload_factor < 80:
            size = "M"
        elif workload_factor < 150:
            size = "L"
        else:
            size = "XL"

        recommended_configs.append({
            'category': 'Warehouse',
            'feature': '',
            'monthly_credits': 0,
            'warehouse_name': f'{other_type.upper()}_{other_name.replace(" ", "_")}',
            'size': size,
            'days_per_month': int(st.session_state.form_data.get(f'other_{i}_days_per_month', 22)),
            'hours_per_day': int(st.session_state.form_data.get(f'other_{i}_hours_per_day', 8)),
            'warehouses_when_active': 1,
            'annual_growth_rate': 20,
            'dev_start_month': 1,
            'go_live_month': 6,
            'ramp_curve': 'Linear',
            'total': 0
        })

    # Add default configuration if no workloads are defined
    if not recommended_configs:
        recommended_configs = [
            {
                'category': 'Warehouse',
                'feature': '',
                'monthly_credits': 0,
                'warehouse_name': 'COMPUTE_WH',
                'size': 'M',
                'days_per_month': 22,
                'hours_per_day': 8,
                'warehouses_when_active': 1,
                'annual_growth_rate': 20,
                'dev_start_month': 1,
                'go_live_month': 6,
                'ramp_curve': 'Linear',
                'total': 0
            }
        ]

    # ONLY add Snowpipe if there are actual pipelines defined AND user specifically wants it
    # Remove the automatic Snowpipe addition - let users add it manually if needed
    # The old code that automatically added Snowpipe has been removed

    return recommended_configs

def display_interactive_warehouse_table():
    """
    Display and manage an interactive warehouse configuration table with optimized sizing
    """
    st.markdown("### 🏭 Warehouse Configuration")
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("🔄 Regenerate Recommendations", key="regen_warehouse_recommendations"):
            st.session_state.warehouse_data = get_recommended_warehouse_config()
            st.rerun()
    with col2:
        if st.button("🗑️ Clear All Warehouses", key="clear_all_warehouses"):
            st.session_state.warehouse_data = []
            st.rerun()

    # Initialize warehouse data in session state if it doesn't exist
    if 'warehouse_data' not in st.session_state:
        # Start with recommended configurations instead of empty list
        st.session_state.warehouse_data = get_recommended_warehouse_config()

    # Define options for dropdowns
    category_options = ["", "Warehouse", "Serverless", "Cloud Services"]
    serverless_features = ["", "Snowpipe", "Auto Clustering", "Materialized Views", "Search Optimization"]
    warehouse_sizes = ["XS", "S", "M", "L", "XL", "2XL", "3XL", "4XL", "5XL", "6XL"]
    ramp_curves = ["Slowest", "Slow", "Linear", "Fast", "Fastest"]

    # Credit rates per warehouse size per hour
    credit_rates = {
        "XS": 1, "S": 2, "M": 4, "L": 8, "XL": 16,
        "2XL": 32, "3XL": 64, "4XL": 128, "5XL": 256, "6XL": 512
    }

    def calculate_warehouse_credits(size, days_per_month, hours_per_day, warehouses_when_active,
                                  annual_growth_rate, dev_start_month, go_live_month, ramp_curve, current_month=12):
        """Calculate total credits for a warehouse configuration"""
        if not size or not days_per_month or not hours_per_day:
            return 0

        base_credits_per_hour = credit_rates.get(size, 0)
        monthly_base_credits = base_credits_per_hour * hours_per_day * days_per_month * warehouses_when_active

        # Apply ramp-up curve from dev start to go live
        if current_month < dev_start_month:
            return 0
        elif current_month < go_live_month:
            progress = (current_month - dev_start_month) / (go_live_month - dev_start_month)
            curve_multipliers = {
                "Slowest": progress ** 3, "Slow": progress ** 2, "Linear": progress,
                "Fast": progress ** 0.5, "Fastest": progress ** 0.25
            }
            ramp_factor = curve_multipliers.get(ramp_curve, progress)
            return monthly_base_credits * ramp_factor
        else:
            months_post_go_live = current_month - go_live_month
            growth_factor = (1 + annual_growth_rate / 100) ** (months_post_go_live / 12)
            return monthly_base_credits * growth_factor

    # Add new row button
    if st.button("➕ Add New Row", key="add_warehouse_row_unique"):
        new_row = {
            'category': '', 'feature': '', 'monthly_credits': 0, 'warehouse_name': '',
            'size': 'XS', 'days_per_month': 22, 'hours_per_day': 8, 'warehouses_when_active': 1,
            'annual_growth_rate': 20, 'dev_start_month': 1, 'go_live_month': 6, 'ramp_curve': 'Linear', 'total': 0
        }
        st.session_state.warehouse_data.append(new_row)
        st.rerun()

    # Column proportions
    col_headers = st.columns([1.6, 2.4, 0.8, 3.0, 1.1, 0.5, 0.6, 0.6, 0.6, 0.6, 0.6, 1.5, 0.9, 0.5])
    headers = ["Category", "Feature", "Credits", "Warehouse Name", "Size",
               "Days", "Hours", "Multi", "Grwth%", "Start", "Live", "Ramp", "Total Monthly", "Del"]

    # Display headers with compact styling
    for i, header in enumerate(headers):
        col_headers[i].markdown(f'<p class="warehouse-table-header">{header}</p>', unsafe_allow_html=True)

    # Display and edit each row
    total_credits = 0
    rows_to_delete = []

    for idx, row in enumerate(st.session_state.warehouse_data):
        cols = st.columns([1.6, 2.4, 0.8, 3.0, 1.1, 0.6, 0.6, 0.6, 0.6, 0.6, 0.6, 1.5, 0.9, 0.5])

        with cols[0]:  # Category
            category = st.selectbox(
                "Cat", category_options,
                index=category_options.index(row['category']) if row['category'] in category_options else 0,
                key=f"wh_category_{idx}", label_visibility="collapsed"
            )
            row['category'] = category

        with cols[1]:  # Feature
            if category == "Serverless":
                feature = st.selectbox(
                    "Feature", serverless_features,
                    index=serverless_features.index(row['feature']) if row['feature'] in serverless_features else 0,
                    key=f"wh_feature_{idx}", label_visibility="collapsed"
                )
                row['feature'] = feature
            else:
                st.write("")
                row['feature'] = ''

        with cols[2]:  # Credits
            if category == "Serverless":
                monthly_credits = st.number_input(
                    "Credits", min_value=0, value=int(row['monthly_credits']),
                    key=f"wh_monthly_credits_{idx}", label_visibility="collapsed"
                )
                row['monthly_credits'] = monthly_credits
                row['total'] = monthly_credits
            else:
                st.write("")
                row['monthly_credits'] = 0

        # For Serverless rows, skip warehouse-specific columns
        if category == "Serverless":
            for i in range(3, 12):
                with cols[i]:
                    st.write("")
            with cols[12]:  # Total
                st.markdown(f"**{row['total']:,.0f}**")
        else:
            # Warehouse-specific columns
            with cols[3]:  # Warehouse Name
                warehouse_name = st.text_input(
                    "Name", value=row['warehouse_name'],
                    key=f"wh_warehouse_name_{idx}", label_visibility="collapsed",
                    placeholder="Warehouse Name"
                )
                row['warehouse_name'] = warehouse_name

            with cols[4]:  # Size
                current_size = row.get('size', 'XS')
                if current_size not in warehouse_sizes:
                    current_size = 'XS'
                    row['size'] = current_size

                size = st.selectbox(
                    "Size", warehouse_sizes,
                    index=warehouse_sizes.index(current_size),
                    key=f"wh_size_{idx}",
                    label_visibility="collapsed"
                )
                row['size'] = size

            with cols[5]:  # Days
                days_per_month = st.number_input(
                    "Days", min_value=1, max_value=31, value=int(row['days_per_month']),
                    key=f"wh_days_per_month_{idx}", label_visibility="collapsed"
                )
                row['days_per_month'] = days_per_month

            with cols[6]:  # Hours
                hours_per_day = st.number_input(
                    "Hours", min_value=1, max_value=24, value=int(row['hours_per_day']),
                    key=f"wh_hours_per_day_{idx}", label_visibility="collapsed"
                )
                row['hours_per_day'] = hours_per_day

            with cols[7]:  # Multi-Cluster
                warehouses_when_active = st.number_input(
                    "Multi", min_value=1, max_value=100, value=int(row['warehouses_when_active']),
                    key=f"wh_warehouses_when_active_{idx}", label_visibility="collapsed"
                )
                row['warehouses_when_active'] = warehouses_when_active

            with cols[8]:  # Growth Rate
                annual_growth_rate = st.number_input(
                    "Growth", min_value=1, max_value=100, value=int(row['annual_growth_rate']),
                    key=f"wh_annual_growth_rate_{idx}", label_visibility="collapsed"
                )
                row['annual_growth_rate'] = annual_growth_rate

            with cols[9]:  # Dev Start
                dev_start_month = st.number_input(
                    "Dev", min_value=1, max_value=60, value=int(row['dev_start_month']),
                    key=f"wh_dev_start_month_{idx}", label_visibility="collapsed"
                )
                row['dev_start_month'] = dev_start_month

            with cols[10]:  # Go Live
                go_live_month = st.number_input(
                    "Live", min_value=1, max_value=60, value=int(row['go_live_month']),
                    key=f"wh_go_live_month_{idx}", label_visibility="collapsed"
                )
                row['go_live_month'] = go_live_month

            with cols[11]:  # Ramp Curve
                current_ramp = row.get('ramp_curve', 'Linear')
                if current_ramp not in ramp_curves:
                    current_ramp = 'Linear'
                    row['ramp_curve'] = current_ramp

                ramp_curve = st.selectbox(
                    "Ramp", ramp_curves,
                    index=ramp_curves.index(current_ramp),
                    key=f"wh_ramp_curve_{idx}",
                    label_visibility="collapsed"
                )
                row['ramp_curve'] = ramp_curve

            # Calculate total credits for warehouse
            if category == "Warehouse":
                calculated_total = calculate_warehouse_credits(
                    size, days_per_month, hours_per_day, warehouses_when_active,
                    annual_growth_rate, dev_start_month, go_live_month, ramp_curve
                )
                row['total'] = calculated_total
            else:
                row['total'] = 0

            with cols[12]:  # Total
                st.markdown(f"<p style='text-align: right; font-weight: bold; margin: 0;'>{row['total']:,.0f}</p>", unsafe_allow_html=True)
            # with cols[12]:  # Total
            #     st.markdown(f"**{row['total']:,.0f}**")

        # Add to total
        total_credits += row['total']

        # Delete button
        with cols[13]:
            if st.button("🗑️", key=f"wh_delete_{idx}", help="Delete this row"):
                rows_to_delete.append(idx)

    # Delete rows (in reverse order to maintain indices)
    for idx in reversed(rows_to_delete):
        del st.session_state.warehouse_data[idx]
        st.rerun()

    # Display totals
    st.markdown("---")

    # Calculate storage costs
    storage_tb = float(st.session_state.form_data.get('final_raw_volume', 0))
    storage_factor = 0.4 * 1.2 * 1.1 * 1.15  # compression * time_travel * fail_safe * cloning
    effective_storage_tb = storage_tb * storage_factor
    storage_price_per_tb = st.session_state.form_data.get('storage_price', 23.0)
    monthly_storage_cost = effective_storage_tb * storage_price_per_tb

    col1, col2, col3, col4 = st.columns([1, 3, 3, 1])
    with col2:
        st.markdown(f"### Total Monthly Credits: <span style='color: #1f77b4; font-weight: bold;'>{total_credits:,.0f}</span>", unsafe_allow_html=True)
        # st.markdown("### Total Monthly Credits:")
        # st.markdown(f"### **{total_credits:,.0f}**")
    with col3:
        st.markdown(f"### Total Storage Costs Per Month: <span style='color: #1f77b4; font-weight: bold;'>${monthly_storage_cost:,.0f}</span>", unsafe_allow_html=True)
        # st.markdown("### Total Storage Costs Per Month:", )
        # st.markdown(f"### **${monthly_storage_cost:,.0f}**")   # Store total in session state for use in other calculations
    st.session_state.form_data['total_monthly_credits'] = total_credits

    return st.session_state.warehouse_data, total_credits

def calculate_consumption_from_interactive_table(warehouse_data, num_years=3):
    """
    Calculate consumption estimates based on the interactive table data
    USES EXACT SAME LOGIC AS THE INTERACTIVE TABLE
    """
    # Define credit rates per warehouse size per hour (SAME AS TABLE)
    credit_rates = {
        "XS": 1, "S": 2, "M": 4, "L": 8, "XL": 16,
        "2XL": 32, "3XL": 64, "4XL": 128, "5XL": 256, "6XL": 512
    }

    def calculate_warehouse_credits(size, days_per_month, hours_per_day, warehouses_when_active,
                                  annual_growth_rate, dev_start_month, go_live_month, ramp_curve, current_month=12):
        """Calculate total credits for a warehouse configuration - EXACT SAME LOGIC AS TABLE"""
        if not size or not days_per_month or not hours_per_day:
            return 0

        base_credits_per_hour = credit_rates.get(size, 0)
        monthly_base_credits = base_credits_per_hour * hours_per_day * days_per_month * warehouses_when_active

        # Apply ramp-up curve from dev start to go live
        if current_month < dev_start_month:
            return 0
        elif current_month < go_live_month:
            progress = (current_month - dev_start_month) / (go_live_month - dev_start_month)
            curve_multipliers = {
                "Slowest": progress ** 3, "Slow": progress ** 2, "Linear": progress,
                "Fast": progress ** 0.5, "Fastest": progress ** 0.25
            }
            ramp_factor = curve_multipliers.get(ramp_curve, progress)
            return monthly_base_credits * ramp_factor
        else:
            months_post_go_live = current_month - go_live_month
            growth_factor = (1 + annual_growth_rate / 100) ** (months_post_go_live / 12)
            return monthly_base_credits * growth_factor

    annual_totals = []

    for year in range(1, num_years + 1):
        year_total = 0
        current_month = year * 12  # End of year calculation

        for row in warehouse_data:
            if row['category'] == 'Serverless':
                # Serverless: fixed monthly credits with growth
                base_monthly = row['monthly_credits']
                if year > 1:
                    growth_factor = (1 + row['annual_growth_rate'] / 100) ** (year - 1)
                    base_monthly *= growth_factor
                year_total += base_monthly * 12
            elif row['category'] == 'Warehouse':
                # Warehouse: calculate using EXACT SAME LOGIC AS TABLE
                monthly_credits = calculate_warehouse_credits(
                    row['size'], row['days_per_month'], row['hours_per_day'],
                    row['warehouses_when_active'], row['annual_growth_rate'],
                    row['dev_start_month'], row['go_live_month'], row['ramp_curve'],
                    current_month
                )
                year_total += monthly_credits * 12  # Convert to annual

        annual_totals.append(year_total)

    return annual_totals

def sync_widget_to_form_data(key):
    """
    A generic callback to sync a widget's state into the form_data dictionary.
    """
    st.session_state.form_data[key] = st.session_state[key]

def generate_pptx(pptx_content, template_file=None):
    """
    Generate a PowerPoint file from pptx_content (list of dicts with 'title', 'body', and optional 'image').
    Returns the path to the generated pptx file.
    """

    # If a template file (as a file-like object) is provided, use it. Otherwise, create a blank presentation.

    prs = Presentation(template_file) if template_file else Presentation()
    for slide_data in pptx_content:
        # Use a standard layout index. Note: This might need adjustment if the template's layout order differs.
        # Layout 5 is typically 'Title Only'. Layout 1 is 'Title and Content'. We'll use 5 for flexibility.
        try:
            slide_layout = prs.slide_layouts[5]
        except IndexError:
            # Fallback to the first layout if the template has fewer than 6 layouts
            st.warning("Template has limited layouts. Using the first available layout.")
            slide_layout = prs.slide_layouts[3]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = slide_data.get('title', '')
        # Add body text
        body_text = slide_data.get('body', '')
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = body_text
        p.font.size = Pt(14) # Slightly smaller for better fit in templates
        # Add image if present
        image_path = slide_data.get('image')
        if image_path and os.path.exists(image_path):
            try:
                # Attempt to place the image. Adjust as needed for your template's design.
                img_left = Inches(1.5)
                img_top = Inches(2.5)
                img_height = Inches(4.5)
                pic = slide.shapes.add_picture(image_path, img_left, img_top, height=img_height)
            except Exception as e:
                st.error(f"Could not add image to PowerPoint slide: {e}")
    # Save to a temporary file
    import tempfile
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.flush()
    return tmp.name

# --- PNG GENERATION SECTION ---

def generate_architecture_diagram_png(form_data, custom_prompt=""):
    """
    Generate a future state architectural diagram in PNG format using matplotlib.
    This is a simplified version of the SVG diagram, but covers the same logical structure.
    """
    # Extract key information from form data
    customer_name = str(form_data.get('customer_name', 'Customer'))
    data_sources_count = int(form_data.get('data_sources_count', 0))
    pipeline_count = int(form_data.get('pipeline_count', 0))
    analytics_count = int(form_data.get('analytics_workload_count', 0))
    tools = form_data.get('tools', [])
    existing_platforms = form_data.get('existing_data_platform', [])
    # Set up the figure
    fig, ax = plt.subplots(figsize=(12, 7))
    ax.axis('off')
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    # Title
    ax.text(7, 7.7, "Snowflake AI Data Cloud - Future State Architecture", fontsize=16, fontweight='bold', ha='center')
    ax.text(7, 7.4, customer_name, fontsize=12, ha='center', color='gray')
    # Data Sources (Left)
    source_start_y = 6
    source_spacing = 1.5
    for i in range(min(max(data_sources_count, 1), 4)):
        y = source_start_y - i * source_spacing
        source_name = str(form_data.get(f'source_name_{i}', f'Data Source {i+1}'))
        source_type = str(form_data.get(f'source_type_{i}', 'Database'))
        # Draw box
        rect = plt.Rectangle((0.6, y-0.6), 2, 1, linewidth=1, edgecolor='gray', facecolor='white')
        ax.add_patch(rect)
        # Draw circle
        circ = plt.Circle((1.5, y+0.1), 0.2, color='#4CAF50')
        ax.add_patch(circ)
        ax.text(1.5, y+0.1, "DB", color='white', fontsize=8, ha='center', va='center')
        ax.text(1.5, y-0.5, source_name[:15], fontsize=9, ha='center')
        ax.text(1.5, y-0.3, source_type[:15], fontsize=7, ha='center', color='gray')
    # ETL/Pipeline (Center-left)
    etl_y = 4
    pipeline_name = str(form_data.get('pipeline_0_name', 'Data Pipeline'))
    rect = plt.Rectangle((3.5, etl_y-0.5), 2, 1.2, linewidth=1, edgecolor='gray', facecolor='white')
    ax.add_patch(rect)
    rect2 = plt.Rectangle((3.7, etl_y-0.4), 1.6, 1, linewidth=0, edgecolor=None, facecolor='#FF9800')
    ax.add_patch(rect2)
    ax.text(4.5, etl_y, "ETL/ELT", color='white', fontsize=10, ha='center', va='center')
    ax.text(4.5, etl_y-0.2, "Processing", color='white', fontsize=8, ha='center', va='center')
    ax.text(4.5, etl_y+0.4, pipeline_name[:20], fontsize=9, ha='center')
    # Snowflake (Center)
    rect = plt.Rectangle((7, etl_y-0.5), 2, 1.2, linewidth=2, edgecolor='#29B5E8', facecolor='white')
    ax.add_patch(rect)
    rect2 = plt.Rectangle((7.2, etl_y-0.4), 1.6, 1, linewidth=0, edgecolor=None, facecolor='#29B5E8')
    ax.add_patch(rect2)
    ax.text(8, etl_y, "SNOWFLAKE", color='white', fontsize=9, fontweight='bold', ha='center', va='center')
    ax.text(8, etl_y-0.2, "Data Cloud", color='white', fontsize=8, ha='center', va='center')
    ax.text(8, etl_y+0.4, "Data Warehouse", fontsize=9, ha='center')
    # Analytics/Output Tools (Right)
    output_tools = []
    if tools and isinstance(tools, list):
        output_tools.extend([str(tool) for tool in tools[:4]])
    default_tools = ['Power BI', 'Tableau', 'ML Models', 'Data Apps']
    while len(output_tools) < 4:
        for default_tool in default_tools:
            if default_tool not in output_tools and len(output_tools) < 4:
                output_tools.append(default_tool)
    output_start_y = 6
    output_spacing = 1.5
    for i, tool in enumerate(output_tools[:4]):
        y = output_start_y - i * output_spacing
        # Color
        if any(keyword in tool.lower() for keyword in ['bi', 'tableau', 'power']):
            tool_color = '#9C27B0'
        elif any(keyword in tool.lower() for keyword in ['ml', 'model', 'ai']):
            tool_color = '#FF5722'
        elif any(keyword in tool.lower() for keyword in ['app', 'web']):
            tool_color = '#4CAF50'
        else:
            tool_color = '#607D8B'
        rect = plt.Rectangle((11, y-0.6), 2, 1.3, linewidth=1, edgecolor='gray', facecolor='white')
        ax.add_patch(rect)
        rect2 = plt.Rectangle((11.2, y-0.5), 1.7, 1.1, linewidth=0, edgecolor=None, facecolor=tool_color)
        ax.add_patch(rect2)
        ax.text(12, y, tool[:12], color='white', fontsize=8, ha='center', va='center')
        ax.text(12, y-0.2, "Analytics", color='white', fontsize=7, ha='center', va='center')
        ax.text(12, y+0.4, tool[:15], fontsize=9, ha='center')
    # Connections
    source_count = min(max(data_sources_count, 1), 4)
    for i in range(source_count):
        y = source_start_y - i * source_spacing
        ax.annotate('', xy=(3.5, etl_y), xytext=(2.6, y), arrowprops=dict(arrowstyle='->', color='#29B5E8', lw=2))
    ax.annotate('', xy=(7.1, etl_y), xytext=(5.5, etl_y), arrowprops=dict(arrowstyle='->', color='#29B5E8', lw=2))
    for i in range(4):
        y = output_start_y - i * output_spacing
        ax.annotate('', xy=(11, y), xytext=(9, etl_y), arrowprops=dict(arrowstyle='->', color='#29B5E8', lw=2))
    # Footer
    if custom_prompt:
        ax.text(0.5, 0.3, f"Requirements: {custom_prompt[:100]}", fontsize=8, color='gray', ha='left')
    current_time = datetime.now().strftime("%Y-%m-%d %H:%M")
    ax.text(0.5, 0.1, f"Sources: {data_sources_count} | Pipelines: {pipeline_count} | Tools: {len(tools)} | Generated: {current_time}", fontsize=8, color='gray', ha='left')
    # Save to BytesIO
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf.read()


def display_architecture_diagram_with_png():
    """
    Display the architecture diagram as SVG and PNG with buttons.
    """
    st.markdown("### 🏗️ Current State Architecture Diagram")
    custom_prompt = st.text_input("Add a custom note to the diagram (optional):", placeholder="e.g., Phase 1 Implementation, POC Architecture, etc.")
    if st.button("🔄 Generate/Refresh SVG Diagram"):
        st.session_state.diagram_generated = True
        st.session_state.diagram_type = 'svg'
    if st.button("🖼️ Generate/Refresh PNG Diagram"):
        st.session_state.diagram_generated = True
        st.session_state.diagram_type = 'png'
    if st.session_state.get('diagram_generated', False):
        if st.session_state.get('diagram_type', 'svg') == 'svg':
            svg_content = generate_architecture_diagram(st.session_state.form_data, custom_prompt)
            st.session_state.current_svg = svg_content
            # st.components.v1.html(svg_content, height=850, width=750)
            st.image(svg_content, width=1200)
            with st.expander("View SVG Code"):
                st.code(svg_content, language='xml')
        else:
            png_bytes = generate_architecture_diagram_png(st.session_state.form_data, custom_prompt)
            st.session_state.current_png = png_bytes
            st.image(png_bytes, caption="Architecture Diagram (PNG)", use_container_width=True)
    if st.button("💾 Save SVG File"):
        if 'current_svg' in st.session_state:
            svg_content = st.session_state.current_svg
            st.download_button(
                label="Download SVG",
                data=svg_content,
                file_name=f"{st.session_state.form_data.get('customer_name', 'customer')}_architecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.svg",
                mime="image/svg+xml"
            )
        else:
            st.warning("Please generate the SVG diagram first")
    if st.button("💾 Save PNG File"):
        if 'current_png' in st.session_state:
            png_bytes = st.session_state.current_png
            st.download_button(
                label="Download PNG",
                data=png_bytes,
                file_name=f"{st.session_state.form_data.get('customer_name', 'customer')}_architecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png",
                mime="image/png"
            )
        else:
            st.warning("Please generate the PNG diagram first")

def generate_architecture_diagram(form_data, custom_prompt=""):
    """
    Generate a future state architectural diagram in SVG format based on collected data
    """
    from datetime import datetime
    import html
    customer_name = html.escape(str(form_data.get('customer_name', 'Customer')))
    data_sources_count = int(form_data.get('data_sources_count', 0))
    pipeline_count = int(form_data.get('pipeline_count', 0))
    analytics_count = int(form_data.get('analytics_workload_count', 0))
    tools = form_data.get('tools', [])
    existing_platforms = form_data.get('existing_data_platform', [])
    width = 1400
    height = 800
    svg_content = f"""<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg" viewBox="0 0 {width} {height}">
    <defs>
        <marker id="arrowhead" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto" markerUnits="strokeWidth">
            <path d="M0,0 L0,6 L9,3 z" fill="#333"/>
        </marker>
        <style>
        <![CDATA[
            .pulse-line {{
                stroke-dasharray: 5,10;
                animation: dash 2s linear infinite;
            }}
            @keyframes dash {{
                to {{
                    stroke-dashoffset: -15;
                }}
            }}
            .box-shadow {{
                filter: drop-shadow(2px 2px 4px rgba(0,0,0,0.1));
            }}
        ]]>
        </style>
    </defs>
    <rect width="{width}" height="{height}" fill="#fafafa"/>
    <text x="{width//2}" y="40" font-size="20" font-weight="bold" text-anchor="middle" font-family="Arial, sans-serif">
        Snowflake AI Data Cloud - Future State Architecture
    </text>
    <text x="{width//2}" y="65" font-size="16" text-anchor="middle" font-family="Arial, sans-serif" fill="#666">
        {customer_name}
    </text>
    """
    source_start_y = 120
    source_spacing = 140
    svg_content += """
    <!-- Data Sources -->"""
    for i in range(min(max(data_sources_count, 1), 4)):
        y_pos = source_start_y + (i * source_spacing)
        source_name = html.escape(str(form_data.get(f'source_name_{i}', f'Data Source {i+1}')))
        source_type = html.escape(str(form_data.get(f'source_type_{i}', 'Database')))
        svg_content += f"""
    <g class="box-shadow">
        <rect x="50" y="{y_pos}" width="150" height="100" fill="#ffffff" stroke="#ddd" stroke-width="1" rx="8"/>
        <circle cx="125" cy="{y_pos + 35}" r="20" fill="#4CAF50"/>
        <text x="125" y="{y_pos + 40}" text-anchor="middle" font-size="10" fill="white" font-family="Arial, sans-serif">DB</text>
        <text x="125" y="{y_pos + 65}" text-anchor="middle" font-size="12" font-family="Arial, sans-serif">{source_name[:15]}</text>
        <text x="125" y="{y_pos + 80}" text-anchor="middle" font-size="10" fill="#666" font-family="Arial, sans-serif">{source_type[:15]}</text>
    </g>"""
    etl_y = 300
    pipeline_name = html.escape(str(form_data.get('pipeline_0_name', 'Data Pipeline')))
    svg_content += f"""
    <g class="box-shadow">
        <rect x="350" y="{etl_y}" width="180" height="100" fill="#ffffff" stroke="#ddd" stroke-width="1" rx="8"/>
        <rect x="370" y="{etl_y + 20}" width="140" height="60" fill="#FF9800" rx="4"/>
        <text x="440" y="{etl_y + 45}" text-anchor="middle" font-size="12" fill="white" font-family="Arial, sans-serif">ETL/ELT</text>
        <text x="440" y="{etl_y + 60}" text-anchor="middle" font-size="10" fill="white" font-family="Arial, sans-serif">Processing</text>
        <text x="440" y="{etl_y + 95}" text-anchor="middle" font-size="11" font-family="Arial, sans-serif">{pipeline_name[:20]}</text>
    </g>"""
    svg_content += f"""
    <g class="box-shadow">
        <rect x="650" y="{etl_y}" width="200" height="100" fill="#ffffff" stroke="#29B5E8" stroke-width="2" rx="8"/>
        <rect x="670" y="{etl_y + 15}" width="160" height="70" fill="#29B5E8" rx="4"/>
        <text x="750" y="{etl_y + 45}" text-anchor="middle" font-size="14" fill="white" font-weight="bold" font-family="Arial, sans-serif">SNOWFLAKE</text>
        <text x="750" y="{etl_y + 60}" text-anchor="middle" font-size="10" fill="white" font-family="Arial, sans-serif">Data Cloud</text>
        <text x="750" y="{etl_y + 95}" text-anchor="middle" font-size="11" font-family="Arial, sans-serif">Data Warehouse</text>
    </g>"""
    output_tools = []
    if tools and isinstance(tools, list):
        output_tools.extend([html.escape(str(tool)) for tool in tools[:4]])
    default_tools = ['Power BI', 'Tableau', 'ML Models', 'Data Apps']
    while len(output_tools) < 4:
        for default_tool in default_tools:
            if default_tool not in output_tools and len(output_tools) < 4:
                output_tools.append(default_tool)
    svg_content += """
    <!-- Analytics Tools -->"""
    output_start_y = 120
    output_spacing = 140
    for i, tool in enumerate(output_tools[:4]):
        y_pos = output_start_y + (i * output_spacing)
        if any(keyword in tool.lower() for keyword in ['bi', 'tableau', 'power']):
            tool_color = '#9C27B0'
        elif any(keyword in tool.lower() for keyword in ['ml', 'model', 'ai']):
            tool_color = '#FF5722'
        elif any(keyword in tool.lower() for keyword in ['app', 'web']):
            tool_color = '#4CAF50'
        else:
            tool_color = '#607D8B'
        svg_content += f"""
    <g class="box-shadow">
        <rect x="1000" y="{y_pos}" width="150" height="100" fill="#ffffff" stroke="#ddd" stroke-width="1" rx="8"/>
        <rect x="1020" y="{y_pos + 20}" width="110" height="60" fill="{tool_color}" rx="4"/>
        <text x="1075" y="{y_pos + 45}" text-anchor="middle" font-size="10" fill="white" font-family="Arial, sans-serif">{tool[:12]}</text>
        <text x="1075" y="{y_pos + 60}" text-anchor="middle" font-size="8" fill="white" font-family="Arial, sans-serif">Analytics</text>
        <text x="1075" y="{y_pos + 95}" text-anchor="middle" font-size="11" font-family="Arial, sans-serif">{tool[:15]}</text>
    </g>"""
    svg_content += """
    <!-- Data Flow Connections -->"""
    source_count = min(max(data_sources_count, 1), 4)
    for i in range(source_count):
        source_y = source_start_y + (i * source_spacing) + 50
        svg_content += f"""
    <line x1="200" y1="{source_y}" x2="350" y2="{etl_y + 50}" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    <line x1="200" y1="{source_y}" x2="350" y2="{etl_y + 50}" stroke="#29B5E8" stroke-width="2" class="pulse-line"/>"""
    svg_content += f"""
    <line x1="530" y1="{etl_y + 50}" x2="650" y2="{etl_y + 50}" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    <line x1="530" y1="{etl_y + 50}" x2="650" y2="{etl_y + 50}" stroke="#29B5E8" stroke-width="2" class="pulse-line"/>"""
    for i in range(len(output_tools[:4])):
        output_y = output_start_y + (i * output_spacing) + 50
        svg_content += f"""
    <line x1="850" y1="{etl_y + 50}" x2="1000" y2="{output_y}" stroke="#666" stroke-width="2" marker-end="url(#arrowhead)"/>
    <line x1="850" y1="{etl_y + 50}" x2="1000" y2="{output_y}" stroke="#29B5E8" stroke-width="2" class="pulse-line"/>"""
    if custom_prompt:
        custom_prompt_escaped = html.escape(str(custom_prompt)[:100])
        svg_content += f"""
    <text x="50" y="{height - 40}" font-size="10" fill="#666" font-family="Arial, sans-serif">Requirements: {custom_prompt_escaped}</text>"""
    current_time = datetime.now().strftime("%Y-%m-%d %H:%M")
    svg_content += f"""
    <text x="50" y="{height - 20}" font-size="10" fill="#666" font-family="Arial, sans-serif">
        Sources: {data_sources_count} | Pipelines: {pipeline_count} | Tools: {len(tools)} | Generated: {current_time}
    </text>
    </svg>"""
    return svg_content

# Helper function to save SVG to file
def save_diagram_to_file(svg_content, filename="architecture_diagram.svg"):
    """
    Save the SVG content to a file
    """
    try:
        with open(filename, 'w', encoding='utf-8') as f:
            f.write(svg_content)
        print(f"Diagram saved to {filename}")
        return True
    except Exception as e:
        print(f"Error saving diagram: {e}")
        return False
# Example usage function
def create_sample_diagram():
    """
    Create a sample diagram with test data
    """
    sample_data = {
        'customer_name': 'Acme Corporation',
        'data_sources_count': 3,
        'pipeline_count': 2,
        'analytics_workload_count': 4,
        'tools': ['Power BI', 'Tableau', 'Python ML', 'Streamlit'],
        'source_name_0': 'CRM Database',
        'source_type_0': 'PostgreSQL',
        'source_name_1': 'ERP System',
        'source_type_1': 'SAP',
        'source_name_2': 'Web Analytics',
        'source_type_2': 'JSON API',
        'pipeline_0_name': 'Daily ETL Process'
    }
    svg_diagram = generate_architecture_diagram(sample_data, "Real-time analytics with ML capabilities")
    save_diagram_to_file(svg_diagram, "sample_architecture.svg")
    return svg_diagram
# Test the function
# This block is guarded to prevent it from running inside Streamlit
if __name__ == "__main__" and not hasattr(st, 'script_runner'):
    # Create and save a sample diagram
    diagram = create_sample_diagram()
    print("Sample diagram created successfully!")
    print(f"SVG length: {len(diagram)} characters")
def display_architecture_diagram():
    """
    Display the architecture diagram with customization options
    """
    st.markdown("### 🏗️ Current State Architecture Diagram")
    # Custom prompt input
    custom_prompt = st.text_input("Add a custom note to the diagram (optional):",
                                  placeholder="e.g., Phase 1 Implementation, POC Architecture, etc.")
    if st.button("🔄 Generate/Refresh Diagram"):
        st.session_state.diagram_generated = True
    if st.button("💾 Save SVG File"):
        if 'current_svg' in st.session_state:
             svg_content = st.session_state.current_svg
             st.download_button(
                 label="Download SVG",
                 data=svg_content,
                 file_name=f"{st.session_state.form_data.get('customer_name', 'customer')}_architecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.svg",
                 mime="image/svg+xml"
                 )
        else:
             st.warning("Please generate the diagram first")
    st.divider()
    if st.session_state.get('diagram_generated', False) or st.button("🎨 Show Architecture Diagram"):
        # Generate the SVG
        svg_content = generate_architecture_diagram(st.session_state.form_data, custom_prompt)
        st.session_state.current_svg = svg_content
        # Display the SVG
        st.components.v1.html(svg_content, height=550)  # adjust height as needed
        # Show SVG code in expander
        with st.expander("View SVG Code"):
            st.code(svg_content, language='xml')
# --- START: MODIFIED CALCULATION AND CHARTING FUNCTIONS ---
def calculate_consumption_estimates(num_years):
    """
    Calculate detailed storage and compute estimates for a multi-year period.
    Returns storage estimates, compute estimates, and annual projections.
    """
    # --- Base Year 1 Calculations ---
    # Storage Factors
    compression_factor = 0.4
    time_travel_factor = 1.2
    fail_safe_factor = 1.1
    cloning_factor = 1.15
    total_storage_factor = compression_factor * time_travel_factor * fail_safe_factor * cloning_factor
    # Storage Year 1
    initial_storage_tb = float(st.session_state.form_data.get('initial_raw_volume', 0))
    final_storage_y1_tb = float(st.session_state.form_data.get('final_raw_volume', 0))
    final_effective_storage_y1 = final_storage_y1_tb * total_storage_factor
    # Compute Year 1
    compute_estimates = []
    # (The existing detailed per-workload calculation logic is used to establish the baseline for Year 1)
    # ETL/ELT Compute
    pipeline_count = int(st.session_state.form_data.get('pipeline_count', 0))
    for i in range(pipeline_count):
        pipeline_name = st.session_state.form_data.get(f'pipeline_{i}_name', f'Pipeline {i+1}')
        runtime_minutes = int(st.session_state.form_data.get(f'pipeline_{i}_runtime', 30))
        jobs_per_day = int(st.session_state.form_data.get(f'pipeline_{i}_jobs_per_day', 24))
        days_per_month = int(st.session_state.form_data.get(f'pipeline_{i}_days_per_month', 22))
        concurrent_jobs = int(st.session_state.form_data.get(f'pipeline_{i}_concurrent_jobs', 5))
        volume_per_job = float(st.session_state.form_data.get(f'pipeline_{i}_volume_per_job', 10.0))
        workload_intensity = volume_per_job * concurrent_jobs
        if workload_intensity < 50: warehouse_size, credits_per_hour = "X-Small", 1
        elif workload_intensity < 200: warehouse_size, credits_per_hour = "Small", 2
        elif workload_intensity < 500: warehouse_size, credits_per_hour = "Medium", 4
        elif workload_intensity < 1000: warehouse_size, credits_per_hour = "Large", 8
        else: warehouse_size, credits_per_hour = "X-Large", 16
        total_runtime_hours = (runtime_minutes / 60) * jobs_per_day * days_per_month
        monthly_credits = total_runtime_hours * credits_per_hour
        multi_cluster = concurrent_jobs > 5
        cluster_count = min(10, max(2, concurrent_jobs // 3)) if multi_cluster else 1
        if multi_cluster: monthly_credits *= (cluster_count * 0.7)
        compute_estimates.append({'name': pipeline_name, 'type': 'ETL/ELT', 'warehouse_size': warehouse_size, 'credits_per_hour': credits_per_hour, 'monthly_credits': monthly_credits, 'cluster_count': cluster_count, 'multi_cluster': multi_cluster})
    # Analytics Compute
    analytics_count = int(st.session_state.form_data.get('analytics_workload_count', 0))
    for i in range(analytics_count):
        analytics_name = st.session_state.form_data.get(f'analytics_{i}_name', f'Analytics {i+1}')
        hours_per_day = int(st.session_state.form_data.get(f'analytics_{i}_hours_per_day', 8))
        days_per_month = int(st.session_state.form_data.get(f'analytics_{i}_days_per_month', 22))
        concurrent_queries = int(st.session_state.form_data.get(f'analytics_{i}_concurrent_queries', 20))
        power_users_pct = int(st.session_state.form_data.get(f'analytics_{i}_power_users', 10))
        caching_pct = int(st.session_state.form_data.get(f'analytics_{i}_caching', 50))
        complexity_factor = concurrent_queries * (1 + power_users_pct/100)
        if complexity_factor < 15: warehouse_size, credits_per_hour = "X-Small", 1
        elif complexity_factor < 40: warehouse_size, credits_per_hour = "Small", 2
        elif complexity_factor < 80: warehouse_size, credits_per_hour = "Medium", 4
        elif complexity_factor < 150: warehouse_size, credits_per_hour = "Large", 8
        else: warehouse_size, credits_per_hour = "X-Large", 16
        base_hours = hours_per_day * days_per_month
        cache_efficiency = 1 - (caching_pct / 100 * 0.3)
        effective_hours = base_hours * cache_efficiency
        monthly_credits = effective_hours * credits_per_hour
        multi_cluster = concurrent_queries > 20
        cluster_count = min(10, max(2, concurrent_queries // 10)) if multi_cluster else 1
        if multi_cluster: monthly_credits *= (cluster_count * 0.8)
        compute_estimates.append({'name': analytics_name, 'type': 'Analytics', 'warehouse_size': warehouse_size, 'credits_per_hour': credits_per_hour, 'monthly_credits': monthly_credits, 'cluster_count': cluster_count, 'multi_cluster': multi_cluster})
    # Other Workloads Compute
    if st.session_state.form_data.get('has_other_workloads', 'No') == 'Yes':
        other_count = int(st.session_state.form_data.get('other_workload_count', 0))
        for i in range(other_count):
            other_name = st.session_state.form_data.get(f'other_{i}_name', f'Other {i+1}')
            other_type = st.session_state.form_data.get(f'other_{i}_type', 'Other')
            hours_per_day = int(st.session_state.form_data.get(f'other_{i}_hours_per_day', 8))
            days_per_month = int(st.session_state.form_data.get(f'other_{i}_days_per_month', 22))
            concurrent_queries = int(st.session_state.form_data.get(f'other_{i}_concurrent_queries', 10))
            data_volume = float(st.session_state.form_data.get(f'other_{i}_data_volume', 5.0))
            workload_factor = concurrent_queries * (data_volume / 10)
            if workload_factor < 10: warehouse_size, credits_per_hour = "X-Small", 1
            elif workload_factor < 30: warehouse_size, credits_per_hour = "Small", 2
            elif workload_factor < 80: warehouse_size, credits_per_hour = "Medium", 4
            elif workload_factor < 150: warehouse_size, credits_per_hour = "Large", 8
            else: warehouse_size, credits_per_hour = "X-Large", 16
            monthly_credits = hours_per_day * days_per_month * (credits_per_hour * 2)
            multi_cluster = concurrent_queries > 15
            cluster_count = min(8, max(2, concurrent_queries // 8)) if multi_cluster else 1
            if multi_cluster: monthly_credits *= (cluster_count * 0.75)
            compute_estimates.append({'name': other_name, 'type': other_type, 'warehouse_size': warehouse_size, 'credits_per_hour': credits_per_hour, 'monthly_credits': monthly_credits, 'cluster_count': cluster_count, 'multi_cluster': multi_cluster})
    total_monthly_credits_y1 = sum([est['monthly_credits'] for est in compute_estimates])
    # --- Multi-Year Projections ---
    annual_growth_rate = 1 + (st.session_state.form_data.get('annual_growth_rate', 20) / 100.0)
    storage_by_year = []
    compute_by_year = []
    # Initialize with Year 1 values
    current_year_storage_tb = final_storage_y1_tb
    current_year_annual_credits = total_monthly_credits_y1 * 12
    for year in range(1, num_years + 1):
        if year > 1:
            current_year_storage_tb *= annual_growth_rate
            current_year_annual_credits *= annual_growth_rate
        storage_by_year.append(current_year_storage_tb * total_storage_factor)
        compute_by_year.append(current_year_annual_credits)
    return {
        'storage': {
            'by_year': storage_by_year,  # List of final effective storage TB per year
        },
        'compute': {
            'estimates': compute_estimates, # Y1 base estimates for warehouse breakdown
            'by_year': compute_by_year, # List of total annual credits per year
            'total_monthly_credits_y1': total_monthly_credits_y1
        },
        'years': list(range(1, num_years + 1)),
    }
def create_consumption_charts(consumption_data):
    """
    Create Streamlit charts for annual storage and compute consumption.
    """
    st.markdown("### 📊 Annual Projections")
    col1, col2 = st.columns(2)
    with col1:
        # Storage Growth Chart (Annual)
        storage_df = pd.DataFrame({
            'Year': consumption_data['years'],
            'Projected Storage (TB)': consumption_data['storage']['by_year']
        })
        fig_storage = px.bar(storage_df, x='Year', y='Projected Storage (TB)',
                             title='Projected Annual Storage (End of Year)',
                             labels={'Projected Storage (TB)': 'Effective Storage (TB)'})
        fig_storage.update_layout(height=400, showlegend=False)
        st.plotly_chart(fig_storage, use_container_width=True)
    with col2:
        # Compute Credits Chart (Annual)
        compute_df = pd.DataFrame({
            'Year': consumption_data['years'],
            'Annual Credits': consumption_data['compute']['by_year']
        })
        fig_compute = px.bar(compute_df, x='Year', y='Annual Credits',
                             title='Projected Annual Compute Credits',
                             color='Annual Credits',
                             color_continuous_scale='Blues')
        fig_compute.update_layout(height=400, showlegend=False)
        st.plotly_chart(fig_compute, use_container_width=True)
    st.markdown("### 🏭 Year 1 Workload Breakdown")
    col3, col4 = st.columns(2)
    warehouse_df = pd.DataFrame(consumption_data['compute']['estimates'])
    with col3:
        # Warehouse Distribution Chart (based on Year 1)
        if not warehouse_df.empty:
            size_counts = warehouse_df['warehouse_size'].value_counts()
            fig_warehouses = px.pie(values=size_counts.values, names=size_counts.index,
                                    title='Distribution of Warehouse Sizes (Year 1)')
            fig_warehouses.update_layout(height=400)
            st.plotly_chart(fig_warehouses, use_container_width=True)
        else:
            st.info("No compute workloads defined.")
    with col4:
        # Credits by Workload Type (based on Year 1)
        if not warehouse_df.empty:
            workload_credits = warehouse_df.groupby('type')['monthly_credits'].sum().reset_index()
            fig_workload = px.pie(workload_credits, values='monthly_credits', names='type',
                                  title='Monthly Credits by Workload Type (Year 1)')
            fig_workload.update_layout(height=400)
            st.plotly_chart(fig_workload, use_container_width=True)
        else:
            st.info("No compute workloads to analyze.")
    return storage_df, compute_df, warehouse_df

def display_cost_estimates(consumption_data, credit_price, storage_price):
    """
    Display detailed cost estimates for the selected projection period inside an expander.
    """
    # Pricing assumptions passed as arguments
    STORAGE_COST_PER_TB_MONTH = storage_price
    CREDIT_COST = credit_price

    # --- Calculations are performed regardless of whether the expander is open ---
    num_years = len(consumption_data['years'])
    final_year_storage_tb = consumption_data['storage']['by_year'][-1]
    final_year_annual_credits = consumption_data['compute']['by_year'][-1]
    final_year_storage_cost = final_year_storage_tb * STORAGE_COST_PER_TB_MONTH * 12
    final_year_compute_cost = final_year_annual_credits * CREDIT_COST
    final_year_total_cost = final_year_storage_cost + final_year_compute_cost
    total_storage_cost = sum(s * STORAGE_COST_PER_TB_MONTH * 12 for s in consumption_data['storage']['by_year'])
    total_compute_cost = sum(c * CREDIT_COST for c in consumption_data['compute']['by_year'])
    total_period_cost = total_storage_cost + total_compute_cost
    cost_breakdown = pd.DataFrame({
        'Cost Type': ['Storage', 'Compute'],
        f'Cost in Year {num_years}': [final_year_storage_cost, final_year_compute_cost]
    })

    # --- All UI elements are placed inside the expander ---
    with st.expander("💰 View Cost Estimates", expanded=False):
        # The markdown header can be removed or kept inside for context
        # st.markdown("### 💰 Cost Estimates")

        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric(f"Annual Cost (Year {num_years})", f"${final_year_total_cost:,.0f}")
        with col2:
            st.metric(f"Total Cost Over {num_years} Years", f"${total_period_cost:,.0f}")
        with col3:
            st.metric("Avg. Annual Cost", f"${(total_period_cost / num_years):,.0f}")

        # Cost breakdown pie chart for the final year
        fig_cost = px.pie(cost_breakdown, values=f'Cost in Year {num_years}', names='Cost Type',
                          title=f'Cost Breakdown in Year {num_years}')
        st.plotly_chart(fig_cost, use_container_width=True)

    # The function still returns the calculated data for use in other parts of the app
    return {
        'final_year_total_cost': final_year_total_cost,
        'total_period_cost': total_period_cost,
        'final_year_storage_cost': final_year_storage_cost,
        'final_year_compute_cost': final_year_compute_cost
    }

# --- END: MODIFIED CALCULATION AND CHARTING FUNCTIONS ---
# --- START: NEW AND MODIFIED SNOWFLAKE FUNCTIONS ---

def get_existing_sizings():
    """
    Fetch all existing sizing records from Snowflake to populate a dropdown.
    Returns a dictionary of {display_name: sizing_id}.
    """
    if not SNOWFLAKE_ENABLED:
        return {"No Snowflake Connection": None}

    try:
        session = get_active_session()
        # Fetch SIZING_ID, CUSTOMER_NAME, and CREATED_TIMESTAMP
        sizings_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.SIZING_MAIN") \
                            .select(col("SIZING_ID"), col("CUSTOMER_NAME"), col("CREATED_TIMESTAMP")) \
                            .sort(col("CREATED_TIMESTAMP").desc()) \
                            .to_pandas()

        if sizings_df.empty:
            return {"No sizings found": None}

        # Create a dictionary for the selectbox
        # Format: "Customer Name - SIZING_ID" -> "SIZING_..."
        sizings_dict = {"-- Select a Sizing to Load --": None} # Add a placeholder
        for _, row in sizings_df.iterrows():
            # Handle potential missing customer names
            customer_name = row['CUSTOMER_NAME'] if pd.notna(row['CUSTOMER_NAME']) and row['CUSTOMER_NAME'] else 'N/A'
            display_name = f"{customer_name} - {row['SIZING_ID']}"
            sizings_dict[display_name] = row['SIZING_ID']

        return sizings_dict

    except Exception as e:
        # Return an error message in the dropdown if something goes wrong
        st.sidebar.error(f"Failed to fetch sizings: {e}")
        return {"Error fetching data": None}

def load_data_from_snowflake(sizing_id):
    """
    Search and populate the form data from a previous sizing ID.
    """
    if not SNOWFLAKE_ENABLED:
        st.error("Snowflake connection is not available.")
        return
    try:
        session = get_active_session()
        st.session_state.form_data = {} # Clear existing form data
        # 1. Fetch from SIZING_MAIN
        main_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.SIZING_MAIN").filter(col("SIZING_ID") == sizing_id).to_pandas()

        if main_df.empty:
            st.error(f"No data found for Sizing ID: {sizing_id}")
            return
        main_data = main_df.iloc[0].to_dict()
        # Populate main fields, converting types as needed
        for key, value in main_data.items():
            key_lower = key.lower()
            if pd.isna(value): continue
            if key_lower in ['existing_data_platform', 'tools', 'geographic_distribution']:
                 st.session_state.form_data[key_lower] = json.loads(value)
            elif key_lower in ['dev_start_date', 'go_live_date']:
                 st.session_state.form_data[key_lower] = datetime.strptime(str(value), '%Y-%m-%d').date()
            elif key_lower in ['data_sources_count', 'expected_warehouses', 'total_users', 'data_retention_period', 'pipeline_count', 'analytics_workload_count', 'annual_growth_rate']:
                 st.session_state.form_data[key_lower] = int(value)
            elif key_lower in ['initial_raw_volume_tb', 'final_raw_volume_tb']:
                 st.session_state.form_data[key.lower().replace('_tb','')] = float(value)
            elif key_lower == 'other_workload_count':
                 st.session_state.form_data[key_lower] = int(value) if value is not None else 0
            else:
                 st.session_state.form_data[key_lower] = value

        # 2. Fetch Data Sources
        sources_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.DATA_SOURCES").filter(col("SIZING_ID") == sizing_id).sort("SOURCE_INDEX").to_pandas()
        for i, row in sources_df.iterrows():
            st.session_state.form_data[f'source_name_{i}'] = row['SOURCE_NAME']
            st.session_state.form_data[f'source_type_{i}'] = row['SOURCE_TYPE']
            st.session_state.form_data[f'current_volume_{i}'] = float(row['CURRENT_VOLUME_TB'])
        # 3. Fetch Pipelines
        pipelines_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.PIPELINES").filter(col("SIZING_ID") == sizing_id).sort("PIPELINE_INDEX").to_pandas()
        for i, row in pipelines_df.iterrows():
            st.session_state.form_data[f'pipeline_{i}_name'] = row['PIPELINE_NAME']
            st.session_state.form_data[f'pipeline_{i}_frequency'] = row['FREQUENCY']
            st.session_state.form_data[f'pipeline_{i}_jobs_per_day'] = int(row['JOBS_PER_DAY'])
            st.session_state.form_data[f'pipeline_{i}_volume_per_job'] = float(row['VOLUME_PER_JOB_GB'])
            st.session_state.form_data[f'pipeline_{i}_runtime'] = int(row['RUNTIME_MINUTES'])
            st.session_state.form_data[f'pipeline_{i}_days_per_month'] = int(row['DAYS_PER_MONTH'])
            st.session_state.form_data[f'pipeline_{i}_concurrent_jobs'] = int(row['CONCURRENT_JOBS'])
            st.session_state.form_data[f'pipeline_{i}_peak_duration'] = int(row['PEAK_DURATION_MINUTES'])
        # 4. Fetch Analytics Workloads
        analytics_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.ANALYTICS_WORKLOADS").filter(col("SIZING_ID") == sizing_id).sort("ANALYTICS_INDEX").to_pandas()
        st.session_state.form_data['analytics_workload_count'] = len(analytics_df)
        for i, row in analytics_df.iterrows():
            st.session_state.form_data[f'analytics_{i}_name'] = row['WORKLOAD_NAME']
            st.session_state.form_data[f'analytics_{i}_tool'] = row['PRIMARY_TOOL']
            st.session_state.form_data[f'analytics_{i}_hours_per_day'] = int(row['HOURS_PER_DAY'])
            st.session_state.form_data[f'analytics_{i}_days_per_month'] = int(row['DAYS_PER_MONTH'])
            st.session_state.form_data[f'analytics_{i}_queries_per_day'] = int(row['QUERIES_PER_DAY'])
            st.session_state.form_data[f'analytics_{i}_basic_users'] = int(row['BASIC_USERS_PCT'])
            st.session_state.form_data[f'analytics_{i}_expert_users'] = int(row['EXPERT_USERS_PCT'])
            st.session_state.form_data[f'analytics_{i}_power_users'] = int(row['POWER_USERS_PCT'])
            st.session_state.form_data[f'analytics_{i}_peak_days'] = int(row['PEAK_DAYS'])
            st.session_state.form_data[f'analytics_{i}_peaks_per_day'] = int(row['PEAKS_PER_DAY'])
            st.session_state.form_data[f'analytics_{i}_peak_duration'] = int(row['PEAK_DURATION_MINUTES'])
            st.session_state.form_data[f'analytics_{i}_concurrent_queries'] = int(row['CONCURRENT_QUERIES'])
            st.session_state.form_data[f'analytics_{i}_caching'] = int(row['CACHING_PCT'])
        # 5. Fetch Other Workloads
        other_df = session.table("SIZING_TOOL.CUSTOMER_SIZING.OTHER_WORKLOADS").filter(col("SIZING_ID") == sizing_id).sort("OTHER_INDEX").to_pandas()
        st.session_state.form_data['other_workload_count'] = len(other_df)
        if len(other_df) > 0:
            st.session_state.form_data['has_other_workloads'] = 'Yes'
        for i, row in other_df.iterrows():
             st.session_state.form_data[f'other_{i}_name'] = row['WORKLOAD_NAME']
             st.session_state.form_data[f'other_{i}_type'] = row['WORKLOAD_TYPE']
             st.session_state.form_data[f'other_{i}_hours_per_day'] = int(row['HOURS_PER_DAY'])
             st.session_state.form_data[f'other_{i}_days_per_month'] = int(row['DAYS_PER_MONTH'])
             st.session_state.form_data[f'other_{i}_queries_per_day'] = int(row['QUERIES_PER_DAY'])
             st.session_state.form_data[f'other_{i}_data_volume'] = float(row['DATA_VOLUME_GB'])
             st.session_state.form_data[f'other_{i}_peak_days'] = int(row['PEAK_DAYS'])
             st.session_state.form_data[f'other_{i}_peaks_per_day'] = int(row['PEAKS_PER_DAY'])
             st.session_state.form_data[f'other_{i}_peak_duration'] = int(row['PEAK_DURATION_MINUTES'])
             st.session_state.form_data[f'other_{i}_concurrent_queries'] = int(row['CONCURRENT_QUERIES'])
        st.success(f"Successfully loaded data for Sizing ID: {sizing_id}")
        st.rerun()
    except Exception as e:
        st.error(f"❌ Error loading from Snowflake: {str(e)}")
        # In case of partial load, clear the form data
        st.session_state.form_data = {}
def update_in_snowflake(sizing_id):
    """
    Deletes the old sizing data and writes the new data for the given sizing_id.
    """
    if not SNOWFLAKE_ENABLED:
        st.error("Snowflake connection is not available.")
        return
    try:
        session = get_active_session()
        st.info(f"Starting update for Sizing ID: {sizing_id}...")
        # List of tables to clear
        tables_to_clear = [
            "SIZING_MAIN", "DATA_SOURCES", "PIPELINES",
            "ANALYTICS_WORKLOADS", "OTHER_WORKLOADS"
        ]
        # Delete existing records
        for table in tables_to_clear:
            full_table_name = f"SIZING_TOOL.CUSTOMER_SIZING.{table}"
            session.sql(f"DELETE FROM {full_table_name} WHERE SIZING_ID = '{sizing_id}'").collect()
        st.info("Old records deleted. Writing updated records...")
        # Write the new data using the existing ID
        write_to_snowflake(sizing_id_to_use=sizing_id)
        st.success(f"✅ Data successfully updated in Snowflake for Sizing ID: {sizing_id}")
    except Exception as e:
        st.error(f"❌ Error updating Snowflake data: {str(e)}")

def write_to_snowflake(sizing_id_to_use=None):
    """
    Write all collected form data to Snowflake database.
    Can be used for creating a new sizing or updating an existing one.
    """
    if not SNOWFLAKE_ENABLED:
        st.error("Snowflake connection is not available.")
        return None
    try:
        session = get_active_session()
        # If an ID is provided, use it. Otherwise, generate a new one.
        if sizing_id_to_use:
            sizing_id = sizing_id_to_use
        else:
            sizing_id = f"SIZING_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        # Store the ID in the session state so the UI can reflect it
        st.session_state.form_data['sizing_id'] = sizing_id
        # Prepare main sizing record
        main_record = {
            'SIZING_ID': sizing_id,
            'CUSTOMER_NAME': st.session_state.form_data.get('customer_name', ''),
            'INDUSTRY': st.session_state.form_data.get('industry', ''),
            'CONTACT_NAME': st.session_state.form_data.get('contact_name', ''),
            'CONTACT_EMAIL': st.session_state.form_data.get('contact_email', ''),
            'COMPANY_SIZE': st.session_state.form_data.get('company_size', ''),
            'EXISTING_DATA_PLATFORM': json.dumps(st.session_state.form_data.get('existing_data_platform', [])),
            'BUSINESS_OWNER': st.session_state.form_data.get('business_owner', ''),
            'TECH_OWNER': st.session_state.form_data.get('tech_owner', ''),
            'DEV_START_DATE': str(st.session_state.form_data.get('dev_start_date', '')),
            'GO_LIVE_DATE': str(st.session_state.form_data.get('go_live_date', '')),
            'SUCCESS_METRICS': st.session_state.form_data.get('success_metrics', ''),
            'ROADBLOCKS': st.session_state.form_data.get('roadblocks', ''),
            'RAMP_UP_CURVE': st.session_state.form_data.get('ramp_up_curve', ''),
            'DATA_SOURCES_COUNT': st.session_state.form_data.get('data_sources_count', 0),
            'INITIAL_RAW_VOLUME_TB': st.session_state.form_data.get('initial_raw_volume', 0.0),
            'FINAL_RAW_VOLUME_TB': st.session_state.form_data.get('final_raw_volume', 0.0),
            'ANNUAL_GROWTH_RATE': st.session_state.form_data.get('annual_growth_rate', 20),
            'TOOLS': json.dumps(st.session_state.form_data.get('tools', [])),
            'DATA_RETENTION_PERIOD': st.session_state.form_data.get('data_retention_period', 12),
            'DATA_SENSITIVITY': st.session_state.form_data.get('data_sensitivity', ''),
            'EXPECTED_WAREHOUSES': st.session_state.form_data.get('expected_warehouses', 0),
            'TOTAL_USERS': st.session_state.form_data.get('total_users', 0),
            'QUERY_COMPLEXITY': st.session_state.form_data.get('query_complexity', ''),
            'HAS_OTHER_WORKLOADS': st.session_state.form_data.get('has_other_workloads', 'No'),
            'OTHER_WORKLOAD_COUNT': st.session_state.form_data.get('other_workload_count', 0),
            'GEOGRAPHIC_DISTRIBUTION': json.dumps(st.session_state.form_data.get('geographic_distribution', [])),
            'HIGH_AVAILABILITY_REQUIREMENTS': st.session_state.form_data.get('high_availability_requirements', ''),
            'DISASTER_RECOVERY_REQUIREMENTS': st.session_state.form_data.get('disaster_recovery_requirements', ''),
            'CREATED_TIMESTAMP': datetime.now()
        }
        main_df = session.create_dataframe([main_record])
        main_df.write.mode("append").save_as_table("SIZING_TOOL.CUSTOMER_SIZING.SIZING_MAIN")
        # Insert data sources
        data_sources_count = int(st.session_state.form_data.get('data_sources_count', 0))
        sources_records = []
        for i in range(data_sources_count):
            sources_records.append({
                'SIZING_ID': sizing_id, 'SOURCE_INDEX': i + 1,
                'SOURCE_NAME': st.session_state.form_data.get(f'source_name_{i}', ''),
                'SOURCE_TYPE': st.session_state.form_data.get(f'source_type_{i}', ''),
                'CURRENT_VOLUME_TB': st.session_state.form_data.get(f'current_volume_{i}', 0.0),
                'CREATED_TIMESTAMP': datetime.now()
            })
        if sources_records:
            source_df = session.create_dataframe(sources_records)
            source_df.write.mode("append").save_as_table("SIZING_TOOL.CUSTOMER_SIZING.DATA_SOURCES")
        # Insert pipeline information
        pipeline_count = int(st.session_state.form_data.get('pipeline_count', 0))
        pipeline_records = []
        for i in range(pipeline_count):
            pipeline_records.append({
                'SIZING_ID': sizing_id, 'PIPELINE_INDEX': i + 1,
                'PIPELINE_NAME': st.session_state.form_data.get(f'pipeline_{i}_name', ''),
                'FREQUENCY': st.session_state.form_data.get(f'pipeline_{i}_frequency', ''),
                'JOBS_PER_DAY': st.session_state.form_data.get(f'pipeline_{i}_jobs_per_day', 0),
                'VOLUME_PER_JOB_GB': st.session_state.form_data.get(f'pipeline_{i}_volume_per_job', 0.0),
                'RUNTIME_MINUTES': st.session_state.form_data.get(f'pipeline_{i}_runtime', 0),
                'DAYS_PER_MONTH': st.session_state.form_data.get(f'pipeline_{i}_days_per_month', 0),
                'CONCURRENT_JOBS': st.session_state.form_data.get(f'pipeline_{i}_concurrent_jobs', 0),
                'PEAK_DURATION_MINUTES': st.session_state.form_data.get(f'pipeline_{i}_peak_duration', 0),
                'CREATED_TIMESTAMP': datetime.now()
            })
        if pipeline_records:
            pipeline_df = session.create_dataframe(pipeline_records)
            pipeline_df.write.mode("append").save_as_table("SIZING_TOOL.CUSTOMER_SIZING.PIPELINES")
        # Insert analytics workloads
        analytics_count = int(st.session_state.form_data.get('analytics_workload_count', 0))
        analytics_records = []
        for i in range(analytics_count):
            analytics_records.append({
                'SIZING_ID': sizing_id, 'ANALYTICS_INDEX': i + 1,
                'WORKLOAD_NAME': st.session_state.form_data.get(f'analytics_{i}_name', ''),
                'PRIMARY_TOOL': st.session_state.form_data.get(f'analytics_{i}_tool', ''),
                'HOURS_PER_DAY': st.session_state.form_data.get(f'analytics_{i}_hours_per_day', 0),
                'DAYS_PER_MONTH': st.session_state.form_data.get(f'analytics_{i}_days_per_month', 0),
                'QUERIES_PER_DAY': st.session_state.form_data.get(f'analytics_{i}_queries_per_day', 0),
                'BASIC_USERS_PCT': st.session_state.form_data.get(f'analytics_{i}_basic_users', 0),
                'EXPERT_USERS_PCT': st.session_state.form_data.get(f'analytics_{i}_expert_users', 0),
                'POWER_USERS_PCT': st.session_state.form_data.get(f'analytics_{i}_power_users', 0),
                'PEAK_DAYS': st.session_state.form_data.get(f'analytics_{i}_peak_days', 0),
                'PEAKS_PER_DAY': st.session_state.form_data.get(f'analytics_{i}_peaks_per_day', 0),
                'PEAK_DURATION_MINUTES': st.session_state.form_data.get(f'analytics_{i}_peak_duration', 0),
                'CONCURRENT_QUERIES': st.session_state.form_data.get(f'analytics_{i}_concurrent_queries', 0),
                'CACHING_PCT': st.session_state.form_data.get(f'analytics_{i}_caching', 0),
                'CREATED_TIMESTAMP': datetime.now()
            })
        if analytics_records:
            analytics_df = session.create_dataframe(analytics_records)
            analytics_df.write.mode("append").save_as_table("SIZING_TOOL.CUSTOMER_SIZING.ANALYTICS_WORKLOADS")
        # Insert other workloads if they exist
        if st.session_state.form_data.get('has_other_workloads', 'No') == 'Yes':
            other_count = int(st.session_state.form_data.get('other_workload_count', 0))
            other_records = []
            for i in range(other_count):
                other_records.append({
                    'SIZING_ID': sizing_id, 'OTHER_INDEX': i + 1,
                    'WORKLOAD_NAME': st.session_state.form_data.get(f'other_{i}_name', ''),
                    'WORKLOAD_TYPE': st.session_state.form_data.get(f'other_{i}_type', ''),
                    'HOURS_PER_DAY': st.session_state.form_data.get(f'other_{i}_hours_per_day', 0),
                    'DAYS_PER_MONTH': st.session_state.form_data.get(f'other_{i}_days_per_month', 0),
                    'QUERIES_PER_DAY': st.session_state.form_data.get(f'other_{i}_queries_per_day', 0),
                    'DATA_VOLUME_GB': st.session_state.form_data.get(f'other_{i}_data_volume', 0.0),
                    'PEAK_DAYS': st.session_state.form_data.get(f'other_{i}_peak_days', 0),
                    'PEAKS_PER_DAY': st.session_state.form_data.get(f'other_{i}_peaks_per_day', 0),
                    'PEAK_DURATION_MINUTES': st.session_state.form_data.get(f'other_{i}_peak_duration', 0),
                    'CONCURRENT_QUERIES': st.session_state.form_data.get(f'other_{i}_concurrent_queries', 0),
                    'CREATED_TIMESTAMP': datetime.now()
                })
            if other_records:
                other_df = session.create_dataframe(other_records)
                other_df.write.mode("append").save_as_table("SIZING_TOOL.CUSTOMER_SIZING.OTHER_WORKLOADS")
        # Don't show success message if this is part of an update
        if not sizing_id_to_use:
            st.success(f"✅ Data successfully saved to Snowflake! Sizing ID: {sizing_id}")
        return sizing_id
    except Exception as e:
        st.error(f"❌ Error saving to Snowflake: {str(e)}")
        return None

# --- END: NEW AND MODIFIED SNOWFLAKE FUNCTIONS ---

st.markdown("<h1 class='main-header'>Snowflake Customer Sizing Tool</h1>", unsafe_allow_html=True)
st.markdown("""
This tool helps accurately predict storage and compute requirements for new Snowflake customers.
Please provide detailed information in each section to generate an accurate sizing estimate.
""")
# Initialize session state for storing form data
if 'form_data' not in st.session_state:
    st.session_state.form_data = {}
if 'current_section' not in st.session_state:
    st.session_state.current_section = 0
if 'show_summary' not in st.session_state:
    st.session_state.show_summary = False
if 'diagram_generated' not in st.session_state:
    st.session_state.diagram_generated = False
# Define sections
sections = [
    "Customer Information",
    "Use Case Timeline",
    "Technical Overview",
    "Data Loading & Transformation",
    "Analytics Workload",
    "Other Workloads",
    "Summary & Recommendations"
]
# Function to navigate between sections
def next_section():
    if st.session_state.current_section < len(sections) - 1:
        st.session_state.current_section += 1
    else:
        st.session_state.show_summary = True

def prev_section():
    if st.session_state.current_section > 0:
        st.session_state.current_section -= 1

# Sidebar navigation
st.sidebar.title("Navigation")
for i, section in enumerate(sections):
    if st.sidebar.button(section, key=f"nav_{i}"):
        st.session_state.current_section = i
        if i == len(sections) - 1:
            st.session_state.show_summary = True
        else:
            st.session_state.show_summary = False

st.sidebar.markdown("---")
st.sidebar.markdown("### Progress")
progress = (st.session_state.current_section) / (len(sections) - 1)
st.sidebar.progress(progress)
# Customer Information Section
if st.session_state.current_section == 0:
    st.markdown("<h2 class='section-header'>Customer Information</h2>", unsafe_allow_html=True)

    with st.expander("Or, Load Existing Sizing Data"):
        # Get the list of existing sizings from Snowflake
        existing_sizings = get_existing_sizings()
        
        # The selected option will be the display name (the key of the dict)
        selected_sizing_display_name = st.selectbox(
            "Select a previously saved sizing",
            options=list(existing_sizings.keys()), # Use the display names as options
            key="load_sizing_selection"
        )

        # --- MODIFICATION START ---
        # Create columns for the buttons for better layout
        col1, col2 = st.columns(2)

        with col1:
            if st.button("Load Sizing Data"):
                # Get the actual SIZING_ID from the dictionary using the selected display name
                load_id = existing_sizings.get(selected_sizing_display_name)
                
                if load_id:
                    with st.spinner(f"Loading data for {load_id}..."):
                        load_data_from_snowflake(load_id)
                else:
                    st.warning("Please select a valid sizing to load.")
        
        with col2:
            # Add a button to clear the session and start over
            if st.button("Clear / Start New"):
                # Loop through and delete all keys in the session state
                for key in list(st.session_state.keys()):
                    del st.session_state[key]
                # Rerun the app to reflect the cleared state
                st.rerun()
        # --- MODIFICATION END ---

    st.markdown("---")

    col1, col2 = st.columns(2)
    with col1:
        st.text_input("Customer Name",
            key='customer_name',
            on_change=sync_widget_to_form_data,
            args=('customer_name',),
            value=st.session_state.form_data.get('customer_name', ''))

        industry_options = ["Retail", "Healthcare", "Financial Services", "Manufacturing", "Technology", "Media & Entertainment", "Other"]
        try:
            industry_index = industry_options.index(st.session_state.form_data.get('industry'))
        except (ValueError, KeyError):
            industry_index = 0
        st.selectbox("Industry",
            options=industry_options,
            key='industry',
            on_change=sync_widget_to_form_data,
            args=('industry',),
            index=industry_index)

    with col2:
        st.text_input("Primary Contact Name",
            key='contact_name',
            on_change=sync_widget_to_form_data,
            args=('contact_name',),
            value=st.session_state.form_data.get('contact_name', ''))

        st.text_input("Contact Email",
            key='contact_email',
            on_change=sync_widget_to_form_data,
            args=('contact_email',),
            value=st.session_state.form_data.get('contact_email', ''))

    size_options = ["Small (< 100 employees)", "Medium (100-1000 employees)", "Large (> 1000 employees)"]
    try:
        size_index = size_options.index(st.session_state.form_data.get('company_size'))
    except (ValueError, KeyError):
        size_index = 1
    st.radio("Company Size",
        options=size_options,
        key='company_size',
        on_change=sync_widget_to_form_data,
        args=('company_size',),
        index=size_index)

    st.multiselect("Existing Data Platform(s)",
        options=["On-premises data warehouse", "Cloud data warehouse", "SQL Server", "Teradata", "Data lake", "Hadoop", "Postgres", "MySQL", "None/Spreadsheets", "Oracle", "Dynamo DB", "Mongo DB", "Azure SQL", "Other"],
        key='existing_data_platform',
        on_change=sync_widget_to_form_data,
        args=('existing_data_platform',),
        default=st.session_state.form_data.get('existing_data_platform', []))

    st.markdown("---")
    col1, col2 = st.columns([1, 1])
    with col2:
        if st.button("Next: Use Case Timeline", key="next_0"):
            next_section()
            st.rerun()

# Use Case Timeline Section
elif st.session_state.current_section == 1:
    st.markdown("<h2 class='section-header'>Use Case Timeline</h2>", unsafe_allow_html=True)
    st.markdown("<p class='category-label'>Deployment Timeline</p>", unsafe_allow_html=True)

    col1, col2 = st.columns(2)
    with col1:
        st.text_input("Business Owner",
            key='business_owner',
            on_change=sync_widget_to_form_data,
            args=('business_owner',),
            value=st.session_state.form_data.get('business_owner', ''))

        st.date_input("Development Start Date",
            key='dev_start_date',
            on_change=sync_widget_to_form_data,
            args=('dev_start_date',),
            value=st.session_state.form_data.get('dev_start_date', date.today() + timedelta(days=30)))

    with col2:
        st.text_input("Technical Owner",
            key='tech_owner',
            on_change=sync_widget_to_form_data,
            args=('tech_owner',),
            value=st.session_state.form_data.get('tech_owner', ''))

        st.date_input("Go-Live Date",
            key='go_live_date',
            on_change=sync_widget_to_form_data,
            args=('go_live_date',),
            value=st.session_state.form_data.get('go_live_date', date.today() + timedelta(days=90)))

    st.text_area("Success Metrics for this Use Case",
        key='success_metrics',
        on_change=sync_widget_to_form_data,
        args=('success_metrics',),
        value=st.session_state.form_data.get('success_metrics', ''))

    st.text_area("Potential Roadblocks or Risks",
        key='roadblocks',
        on_change=sync_widget_to_form_data,
        args=('roadblocks',),
        value=st.session_state.form_data.get('roadblocks', ''))

    ramp_options = ["Phased", "Gradual", "Very Slow", "Slow", "Moderate", "Fast", "Very Fast", "Aggressive"]
    try:
        ramp_value = st.session_state.form_data.get('ramp_up_curve')
    except (ValueError, KeyError):
        ramp_value = "Moderate"
    st.select_slider("Expected Ramp-up Curve",
        options=ramp_options,
        key='ramp_up_curve',
        on_change=sync_widget_to_form_data,
        args=('ramp_up_curve',),
        value=ramp_value)

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous", key="prev_1"):
            prev_section()
            st.rerun()
    with col3:
        if st.button("Next: Technical Overview", key="next_1"):
            next_section()
            st.rerun()

# Technical Overview Section
elif st.session_state.current_section == 2:
    st.markdown("<h2 class='section-header'>Technical Overview</h2>", unsafe_allow_html=True)

    st.markdown("<p class='category-label'>Database/Data Lake</p>", unsafe_allow_html=True)

    st.number_input(
        "How many relevant data sources for this use case?",
        min_value=0,
        max_value=100,
        key='data_sources_count',
        on_change=sync_widget_to_form_data,
        args=('data_sources_count',),
        value=int(st.session_state.form_data.get('data_sources_count', 1))
    )

    # Create dynamic fields for each data source
    data_sources = []
    for i in range(int(st.session_state.form_data.get('data_sources_count', 0))):
        st.markdown(f"#### Data Source {i+1}")
        col1, col2, col3 = st.columns(3)
        with col1:
            st.text_input(
                f"Source Name",
                key=f"source_name_{i}",
                on_change=sync_widget_to_form_data,
                args=(f"source_name_{i}",),
                value=st.session_state.form_data.get(f'source_name_{i}', f'Source {i+1}')
            )
        with col2:
            source_type_options = ["Relational Database", "SQL Server",  "Oracle", "NoSQL Database", "Bloomberg API", "Excel",
                                   "SAP", "Teradata", "Files (CSV, JSON, etc.)", "API", "Streaming", "R", "Mongo DB", "Kafka",
                                   "Spark", "Python","Cassandra","Hadoop", "Looker", "MicroStrategy","Redis", "Spotfire", "Qlik",
                                   "SAS", "Netezza", "Guidewire", "Dynamo DB", "MySQL", "Salesforce", "Elasticsearch", "Postgres", "S3", "Other"]
            st.selectbox(
                f"Source Type",
                options=source_type_options,
                key=f"source_type_{i}",
                on_change=sync_widget_to_form_data,
                args=(f"source_type_{i}",),
                index=source_type_options.index(st.session_state.form_data.get(f'source_type_{i}', "Relational Database"))
            )
        with col3:
            st.number_input(
                f"Current Raw Data Volume (TB)",
                min_value=0.0,
                key=f"current_volume_{i}",
                on_change=sync_widget_to_form_data,
                args=(f"current_volume_{i}",),
                value=float(st.session_state.form_data.get(f'current_volume_{i}', 1.0))
            )

        data_sources.append({
            'name': st.session_state.form_data.get(f'source_name_{i}'),
            'type': st.session_state.form_data.get(f'source_type_{i}'),
            'volume': st.session_state.form_data.get(f'current_volume_{i}', 0.0)
        })

    st.markdown("---")

    col1, col2, col3 = st.columns(3)
    with col1:
        st.number_input(
            "Initial raw data volume expected in Snowflake (TB)",
            min_value=0.0,
            key='initial_raw_volume',
            on_change=sync_widget_to_form_data,
            args=('initial_raw_volume',),
            value=float(st.session_state.form_data.get('initial_raw_volume', sum([src.get('volume', 0.0) for src in data_sources])))
        )

    with col2:
        st.number_input(
            "Final raw data volume expected after 12 months in Snowflake (TB)",
            min_value=0.0,
            key='final_raw_volume',
            on_change=sync_widget_to_form_data,
            args=('final_raw_volume',),
            value=float(st.session_state.form_data.get('final_raw_volume', st.session_state.form_data.get('initial_raw_volume', 0.0) * 2.0))
        )

    with col3:
        st.number_input(
            "Annual Growth Rate (%) after Year 1",
            min_value=0,
            max_value=100,
            key='annual_growth_rate',
            on_change=sync_widget_to_form_data,
            args=('annual_growth_rate',),
            value=int(st.session_state.form_data.get('annual_growth_rate', 20)),
            help="The compound annual growth rate for storage and compute in years 2-5."
        )

    st.markdown("<p class='category-label'>Use Case Scope</p>", unsafe_allow_html=True)

    st.multiselect(
        "Which applications/tools will be used to load, transform and analyze this data?",
        options=["Snowsight", "Tableau", "Power BI", "Looker", "MicroStrategy", "Qlik", "Informatica",
                 "Talend", "Matillion", "Fivetran", "dbt", "Python", "R", "Spark", "SAP", "Custom Applications",
                 "Excel", "Jupyter", "SAS", "Other"],
        key='tools',
        on_change=sync_widget_to_form_data,
        args=('tools',),
        default=st.session_state.form_data.get('tools', ["Snowsight"])
    )

    st.markdown("### Additional Information")

    col1, col2 = st.columns(2)
    with col1:
        st.slider(
            "Data Retention Period (months)",
            min_value=1,
            max_value=84,
            key='data_retention_period',
            on_change=sync_widget_to_form_data,
            args=('data_retention_period',),
            value=int(st.session_state.form_data.get('data_retention_period', 12))
        )

    with col2:
        sensitivity_options = ["Public", "Internal", "Confidential", "Restricted", "Low", "Medium", "High", "Critical"]
        st.selectbox(
            "Data Sensitivity Level",
            options=sensitivity_options,
            key='data_sensitivity',
            on_change=sync_widget_to_form_data,
            args=('data_sensitivity',),
            index=sensitivity_options.index(st.session_state.form_data.get('data_sensitivity', "Internal"))
        )

    st.number_input(
        "Expected number of warehouses needed",
        min_value=1,
        max_value=20,
        key='expected_warehouses',
        on_change=sync_widget_to_form_data,
        args=('expected_warehouses',),
        value=int(st.session_state.form_data.get('expected_warehouses', 3))
    )

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous", key="prev_2"):
            prev_section()
            st.rerun()
    with col3:
        if st.button("Next: Data Loading & Transformation", key="next_2"):
            next_section()
            st.rerun()

# Data Loading & Transformation Section
elif st.session_state.current_section == 3:
    st.markdown("<h2 class='section-header'>Data Loading & Transformation</h2>", unsafe_allow_html=True)

    # Initialize pipeline count if not set but data exists
    if 'pipeline_count' not in st.session_state.form_data or st.session_state.form_data.get('pipeline_count', 0) == 0:
        # Check if pipeline data exists in form_data and set count accordingly
        pipeline_count = 0
        for key in st.session_state.form_data.keys():
            if key.startswith('pipeline_') and key.endswith('_name'):
                pipeline_index = int(key.split('_')[1])
                pipeline_count = max(pipeline_count, pipeline_index + 1)
        if pipeline_count > 0:
            st.session_state.form_data['pipeline_count'] = pipeline_count

    # Get number of ETL/ELT pipelines
    st.number_input(
        "How many distinct ETL/ELT pipelines will you have?",
        min_value=0,
        max_value=20,
        key='pipeline_count',
        on_change=sync_widget_to_form_data,
        args=('pipeline_count',),
        value=int(st.session_state.form_data.get('pipeline_count', 1))
    )

    # For each pipeline, collect details
    for i in range(int(st.session_state.form_data.get('pipeline_count', 0))):
        st.markdown(f"### Pipeline {i+1}")

        col1, col2 = st.columns(2)
        with col1:
            st.text_input(
                "Pipeline Name",
                key=f"pipeline_{i}_name",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_name",),
                value=st.session_state.form_data.get(f'pipeline_{i}_name', f"Pipeline {i+1}")
            )

            frequency_options = ["Real-time/Streaming", "Every few minutes", "Every 15 min", "Every 30 min", "Hourly", "Every 4 hours", "Daily", "Weekly", "Monthly", "Ad-hoc", "Real-time"]
            st.selectbox(
                "How often does data need to be loaded?",
                options=frequency_options,
                key=f"pipeline_{i}_frequency",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_frequency",),
                index=frequency_options.index(st.session_state.form_data.get(f'pipeline_{i}_frequency', "Daily"))
            )

            st.number_input(
                "Average number of jobs per day",
                min_value=1,
                max_value=5000,
                key=f"pipeline_{i}_jobs_per_day",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_jobs_per_day",),
                value=int(st.session_state.form_data.get(f'pipeline_{i}_jobs_per_day', 24))
            )

            st.number_input(
                "Average data volume processed per job (GB)",
                min_value=0.1,
                max_value=10000.0,
                key=f"pipeline_{i}_volume_per_job",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_volume_per_job",),
                value=float(st.session_state.form_data.get(f'pipeline_{i}_volume_per_job', 10.0))
            )

        with col2:
            st.number_input(
                "Average run time per job (minutes)",
                min_value=1,
                max_value=1440,
                key=f"pipeline_{i}_runtime",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_runtime",),
                value=int(st.session_state.form_data.get(f'pipeline_{i}_runtime', 30))
            )

            st.slider(
                "Average number of days per month you process data",
                min_value=1,
                max_value=31,
                key=f"pipeline_{i}_days_per_month",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_days_per_month",),
                value=int(st.session_state.form_data.get(f'pipeline_{i}_days_per_month', 22))
            )

            st.number_input(
                "Average number of jobs running concurrently at peak",
                min_value=1,
                max_value=100,
                key=f"pipeline_{i}_concurrent_jobs",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_concurrent_jobs",),
                value=int(st.session_state.form_data.get(f'pipeline_{i}_concurrent_jobs', 5))
            )

            st.number_input(
                "Average duration of a peak (minutes)",
                min_value=1,
                max_value=1440,
                key=f"pipeline_{i}_peak_duration",
                on_change=sync_widget_to_form_data,
                args=(f"pipeline_{i}_peak_duration",),
                value=int(st.session_state.form_data.get(f'pipeline_{i}_peak_duration', 60))
            )

    # Additional questions
    st.markdown("### Additional Pipeline Information")

    st.select_slider(
        "Data Transformation Complexity",
        options=["Very Simple", "Simple", "Moderate", "Complex", "Very Complex"],
        key='data_transformation_complexity',
        on_change=sync_widget_to_form_data,
        args=('data_transformation_complexity',),
        value=st.session_state.form_data.get('data_transformation_complexity', "Moderate")
    )

    st.multiselect(
        "Tools used for data pipelines",
        options=["Snowsight", "Informatica", "Matillion", "Fivetran", "dbt", "Python", "R",
                 "Spark", "Tableau", "Looker", "MicroStrategy", "Qlik",  "SAS", "Custom Applications",
                 "Excel", "Jupyter", "Snowflake Tasks", "Snowpipe", "Talend", "Custom Scripts", "Other"],
        key='pipeline_tools',
        on_change=sync_widget_to_form_data,
        args=('pipeline_tools',),
        default=st.session_state.form_data.get('pipeline_tools', ["Snowflake Tasks"])
    )

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous", key="prev_3"):
            prev_section()
            st.rerun()
    with col3:
        if st.button("Next: Analytics Workload", key="next_3"):
            next_section()
            st.rerun()

# Analytics Workload Section
elif st.session_state.current_section == 4:
    st.markdown("<h2 class='section-header'>Analytics Workload</h2>", unsafe_allow_html=True)

    st.markdown("<p class='category-label'>Analytics</p>", unsafe_allow_html=True)

    # Get number of analytics workloads
    st.number_input(
        "How many distinct analytics workloads will you have?",
        min_value=0,
        max_value=10,
        key='analytics_workload_count',
        on_change=sync_widget_to_form_data,
        args=('analytics_workload_count',),
        value=int(st.session_state.form_data.get('analytics_workload_count', 1))
    )

    # For each analytics workload, collect details
    for i in range(int(st.session_state.form_data.get('analytics_workload_count', 0))):
        st.markdown(f"### Analytics Workload {i+1}")

        col1, col2 = st.columns(2)
        with col1:
            st.text_input(
                "Workload Name",
                key=f"analytics_{i}_name",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_name",),
                value=st.session_state.form_data.get(f'analytics_{i}_name', f"Analytics {i+1}")
            )

            tool_options = ["Snowsight", "Tableau", "Power BI", "Looker", "MicroStrategy", "SAP", "Qlik", "Excel", "Spotfire", "Kafka", "SAS", "Python", "Jupyter", "R", "Other"]
            st.selectbox(
                "Primary Analytics Tool",
                options=tool_options,
                key=f"analytics_{i}_tool",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_tool",),
                index=tool_options.index(st.session_state.form_data.get(f'analytics_{i}_tool', "Snowsight"))
            )

            st.slider(
                "Average number of hours per day users use the tool",
                min_value=1,
                max_value=24,
                key=f"analytics_{i}_hours_per_day",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_hours_per_day",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_hours_per_day', 8))
            )

            st.slider(
                "Average number of days per month users use the tool",
                min_value=1,
                max_value=31,
                key=f"analytics_{i}_days_per_month",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_days_per_month",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_days_per_month', 22))
            )

            st.number_input(
                "Average number of queries triggered per day",
                min_value=1,
                max_value=100000,
                key=f"analytics_{i}_queries_per_day",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_queries_per_day",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_queries_per_day', 1000))
            )

        with col2:
            st.markdown("User Distribution:")
            col_a, col_b, col_c = st.columns(3)
            with col_a:
                st.slider(
                    "Basic Users (%)",
                    min_value=0, max_value=100,
                    key=f"analytics_{i}_basic_users",
                    on_change=sync_widget_to_form_data,
                    args=(f"analytics_{i}_basic_users",),
                    value=int(st.session_state.form_data.get(f'analytics_{i}_basic_users', 60))
                )
            with col_b:
                st.slider(
                    "Expert Users (%)",
                    min_value=0, max_value=100,
                    key=f"analytics_{i}_expert_users",
                    on_change=sync_widget_to_form_data,
                    args=(f"analytics_{i}_expert_users",),
                    value=int(st.session_state.form_data.get(f'analytics_{i}_expert_users', 30))
                )
            with col_c:
                st.slider(
                    "Power Users (%)",
                    min_value=0, max_value=100,
                    key=f"analytics_{i}_power_users",
                    on_change=sync_widget_to_form_data,
                    args=(f"analytics_{i}_power_users",),
                    value=int(st.session_state.form_data.get(f'analytics_{i}_power_users', 10))
                )

            total = (st.session_state.form_data.get(f'analytics_{i}_basic_users', 0) +
                     st.session_state.form_data.get(f'analytics_{i}_expert_users', 0) +
                     st.session_state.form_data.get(f'analytics_{i}_power_users', 0))
            if total != 100:
                st.warning(f"User percentages should add up to 100%. Current total: {total}%")

            st.number_input(
                "Average number of peak days in a month",
                min_value=1, max_value=31,
                key=f"analytics_{i}_peak_days",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_peak_days",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_peak_days', 5))
            )

            st.number_input(
                "Average number of peaks per day",
                min_value=1, max_value=24,
                key=f"analytics_{i}_peaks_per_day",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_peaks_per_day",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_peaks_per_day', 2))
            )

            st.number_input(
                "Average duration of a peak (minutes)",
                min_value=1, max_value=1440,
                key=f"analytics_{i}_peak_duration",
                on_change=sync_widget_to_form_data,
                args=(f"analytics_{i}_peak_duration",),
                value=int(st.session_state.form_data.get(f'analytics_{i}_peak_duration', 60))
            )

    col1_queries, col2_caching = st.columns(2)
    with col1_queries:
        st.number_input(
            "Average number of queries running concurrently at peak",
            min_value=1, max_value=1000,
            key=f"analytics_concurrent_queries", # Using a general key as it seems to apply to the whole section now
            on_change=sync_widget_to_form_data,
            args=(f"analytics_concurrent_queries",),
            value=int(st.session_state.form_data.get(f'analytics_concurrent_queries', 20))
        )
    with col2_caching:
        st.slider(
            "Percentage of dashboards leveraging caching/in-memory structures",
            min_value=0, max_value=100,
            key=f"analytics_caching", # Using a general key
            on_change=sync_widget_to_form_data,
            args=(f"analytics_caching",),
            value=int(st.session_state.form_data.get(f'analytics_caching', 50))
        )


    # Additional questions
    st.markdown("### Additional Analytics Information")

    st.number_input(
        "Total number of users who will access Snowflake",
        min_value=1,
        max_value=10000,
        key='total_users',
        on_change=sync_widget_to_form_data,
        args=('total_users',),
        value=int(st.session_state.form_data.get('total_users', 50))
    )

    st.select_slider(
        "Overall Query Complexity",
        options=["Very Simple", "Simple", "Medium", "Moderate", "Complex", "Very Complex"],
        key='query_complexity',
        on_change=sync_widget_to_form_data,
        args=('query_complexity',),
        value=st.session_state.form_data.get('query_complexity', "Moderate")
    )

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous", key="prev_4"):
            prev_section()
            st.rerun()
    with col3:
        if st.button("Next: Other Workloads", key="next_4"):
            next_section()
            st.rerun()

# Other Workloads Section
elif st.session_state.current_section == 5:
    st.markdown("<h2 class='section-header'>Other Workloads</h2>", unsafe_allow_html=True)

    st.markdown("<p class='category-label'>Other Workloads</p>", unsafe_allow_html=True)

    # Ask if there are other workloads
    st.radio(
        "Do you have other workloads to consider?",
        options=["Yes", "No"],
        key='has_other_workloads',
        on_change=sync_widget_to_form_data,
        args=('has_other_workloads',),
        index=["Yes", "No"].index(st.session_state.form_data.get('has_other_workloads', "No"))
    )

    if st.session_state.form_data.get('has_other_workloads') == "Yes":
        # Get number of other workloads
        st.number_input(
            "How many other workloads will you have?",
            min_value=0,
            max_value=10,
            key='other_workload_count',
            on_change=sync_widget_to_form_data,
            args=('other_workload_count',),
            value=int(st.session_state.form_data.get('other_workload_count', 1))
        )

        # For each other workload, collect details
        for i in range(int(st.session_state.form_data.get('other_workload_count', 0))):
            st.markdown(f"### Other Workload {i+1}")

            col1, col2 = st.columns(2)
            with col1:
                st.text_input(
                    "Workload Name",
                    key=f"other_{i}_name",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_name",),
                    value=st.session_state.form_data.get(f'other_{i}_name', f"Other Workload {i+1}")
                )

                other_type_options = ["Data Science/ML", "ML Operations", "Compute Intensive", "Application Backend", "Data Sharing", "Data Governance", "Reporting",
                                      "Data API", "Data Integration", "Stream Processing", "Batch Processing", "Data Sync", "Data Quality", "Micro Services", "Security", "Other"]
                st.selectbox(
                    "Workload Type",
                    options=other_type_options,
                    key=f"other_{i}_type",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_type",),
                    index=other_type_options.index(st.session_state.form_data.get(f'other_{i}_type', "Data Science/ML"))
                )

                st.slider(
                    "Average number of hours per day users use the tool",
                    min_value=1, max_value=24,
                    key=f"other_{i}_hours_per_day",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_hours_per_day",),
                    value=int(st.session_state.form_data.get(f'other_{i}_hours_per_day', 8))
                )

                st.slider(
                    "Average number of days per month users use the tool",
                    min_value=1, max_value=31,
                    key=f"other_{i}_days_per_month",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_days_per_month",),
                    value=int(st.session_state.form_data.get(f'other_{i}_days_per_month', 22))
                )

            with col2:
                st.number_input(
                    "Average number of queries/jobs per day",
                    min_value=1, max_value=100000,
                    key=f"other_{i}_queries_per_day",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_queries_per_day",),
                    value=int(st.session_state.form_data.get(f'other_{i}_queries_per_day', 500))
                )

                st.number_input(
                    "Average data volume processed per query (GB)",
                    min_value=0.1, max_value=10000.0,
                    key=f"other_{i}_data_volume",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_data_volume",),
                    value=float(st.session_state.form_data.get(f'other_{i}_data_volume', 5.0))
                )

                st.number_input(
                    "Average number of peak days in a month",
                    min_value=1, max_value=31,
                    key=f"other_{i}_peak_days",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_peak_days",),
                    value=int(st.session_state.form_data.get(f'other_{i}_peak_days', 5))
                )

                st.number_input(
                    "Average number of peaks per day",
                    min_value=1, max_value=24,
                    key=f"other_{i}_peaks_per_day",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_peaks_per_day",),
                    value=int(st.session_state.form_data.get(f'other_{i}_peaks_per_day', 2))
                )

            col1_peak, col2_concurrent = st.columns(2)
            with col1_peak:
                st.number_input(
                    "Average duration of a peak (minutes)",
                    min_value=1, max_value=1440,
                    key=f"other_{i}_peak_duration",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_peak_duration",),
                    value=int(st.session_state.form_data.get(f'other_{i}_peak_duration', 60))
                )

            with col2_concurrent:
                st.number_input(
                    "Average number of queries running concurrently at peak",
                    min_value=1, max_value=1000,
                    key=f"other_{i}_concurrent_queries",
                    on_change=sync_widget_to_form_data,
                    args=(f"other_{i}_concurrent_queries",),
                    value=int(st.session_state.form_data.get(f'other_{i}_concurrent_queries', 10))
                )

    # Additional questions that might be helpful
    st.markdown("### Additional Considerations")

    st.multiselect(
        "Geographic regions where users will access Snowflake",
        options=["US", "Canada", "North America", "South America", "Europe", "Asia Pacific", "Middle East", "Africa", "Australia/Oceania", "Global"],
        key='geographic_distribution',
        on_change=sync_widget_to_form_data,
        args=('geographic_distribution',),
        default=st.session_state.form_data.get('geographic_distribution', ["North America"])
    )

    st.select_slider(
        "High Availability Requirements",
        options=["Low", "Standard", "Medium", "High", "Very High", "Critical", "Mission Critical"],
        key='high_availability_requirements',
        on_change=sync_widget_to_form_data,
        args=('high_availability_requirements',),
        value=st.session_state.form_data.get('high_availability_requirements', "Standard")
    )

    st.select_slider(
        "Disaster Recovery Requirements",
        options=["Standard", "High", "Very High", "Mission Critical", "Required"],
        key='disaster_recovery_requirements',
        on_change=sync_widget_to_form_data,
        args=('disaster_recovery_requirements',),
        value=st.session_state.form_data.get('disaster_recovery_requirements', "Standard")
    )

    st.markdown("---")
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous", key="prev_5"):
            prev_section()
            st.rerun()
    with col3:
        if st.button("Generate Summary & Recommendations", key="next_5"):
            next_section()
            st.session_state.show_summary = True
            st.rerun()

# Summary & Recommendations Section
elif st.session_state.current_section == 6 or st.session_state.show_summary:
    st.markdown("<h2 class='section-header'>Summary & Recommendations</h2>", unsafe_allow_html=True)

    # --- START: UI ELEMENTS FOR ADJUSTMENTS ---
    st.markdown("#### Adjust Cost and Time Assumptions")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.selectbox(
            "Projection Period (Years)",
            options=[1, 2, 3, 4, 5],
            index=2, # Default to 3 years
            key='projection_years',
            on_change=sync_widget_to_form_data,
            args=('projection_years',)
        )
    with col2:
        st.number_input(
            "Credit Price ($)",
            min_value=0.01,
            max_value=6.00,
            value=3.00,
            step=0.01,
            key='credit_price',
            on_change=sync_widget_to_form_data,
            args=('credit_price',),
            help="Set the price per Snowflake credit."
        )
    with col3:
        st.number_input(
            "Storage Price ($ per TB/month)",
            min_value=1.0,
            max_value=100.0,
            value=23.0,
            step=1.0,
            key='storage_price',
            on_change=sync_widget_to_form_data,
            args=('storage_price',),
            help="Set the price per TB of compressed storage per month."
        )
    st.markdown("---")
    # --- END: UI ELEMENTS ---

    # Get the values after they have been set in the session state by the callback
    num_years = st.session_state.form_data.get('projection_years', 3)
    credit_price = st.session_state.form_data.get('credit_price', 3.00)
    storage_price = st.session_state.form_data.get('storage_price', 23.0)

    # Calculate consumption estimates
    consumption_data = calculate_consumption_estimates(num_years)

    # --- WAREHOUSE CONFIGURATION TABLE ---
    warehouse_data, total_monthly_credits = display_interactive_warehouse_table()

    # Calculate consumption using interactive table data
    interactive_annual_totals = []  # Initialize the variable
    if warehouse_data:
        interactive_annual_totals = calculate_consumption_from_interactive_table(
            warehouse_data,
            st.session_state.form_data.get('projection_years', 3)
        )

    # Update consumption_data structure to use table calculations
    if interactive_annual_totals:
        consumption_data['compute']['by_year'] = interactive_annual_totals
        consumption_data['compute']['total_monthly_credits_y1'] = total_monthly_credits

    # --- COST ESTIMATES USING TABLE DATA ---
    cost_data = display_cost_estimates(consumption_data, credit_price, storage_price)

    # Create charts using the interactive table data
    def create_consumption_charts_from_table(consumption_data, warehouse_data, num_years):
        """
        Create charts using the interactive table data instead of form-based calculations
        """
        st.markdown("### 📊 Annual Projections")
        col1, col2 = st.columns(2)

        # Calculate annual credits from interactive table using SAME LOGIC
        if warehouse_data:
            interactive_annual_totals = calculate_consumption_from_interactive_table(warehouse_data, num_years)
        else:
            interactive_annual_totals = [0] * num_years

        with col1:
            # Storage Growth Chart (keep existing storage calculation)
            storage_df = pd.DataFrame({
                'Year': consumption_data['years'],
                'Projected Storage (TB)': consumption_data['storage']['by_year']
            })
            fig_storage = px.bar(storage_df, x='Year', y='Projected Storage (TB)',
                                 title='Projected Annual Storage (End of Year)',
                                 labels={'Projected Storage (TB)': 'Effective Storage (TB)'})
            fig_storage.update_layout(height=400, showlegend=False)
            st.plotly_chart(fig_storage, use_container_width=True)

        with col2:
            # Compute Credits Chart (use interactive table data)
            compute_df = pd.DataFrame({
                'Year': list(range(1, num_years + 1)),
                'Annual Credits': interactive_annual_totals
            })
            fig_compute = px.bar(compute_df, x='Year', y='Annual Credits',
                                 title='Projected Annual Compute Credits',
                                 color='Annual Credits',
                                 color_continuous_scale='Blues')
            fig_compute.update_layout(height=400, showlegend=False)
            st.plotly_chart(fig_compute, use_container_width=True)

        st.markdown("### 🏭 Year 1 Workload Breakdown")
        col3, col4 = st.columns(2)

        if warehouse_data:
            warehouse_df = pd.DataFrame(warehouse_data)

            with col3:
                # Warehouse Distribution Chart
                if not warehouse_df.empty and 'size' in warehouse_df.columns:
                    warehouse_only = warehouse_df[warehouse_df['category'] == 'Warehouse']
                    if not warehouse_only.empty:
                        size_counts = warehouse_only['size'].value_counts()
                        fig_warehouses = px.pie(values=size_counts.values, names=size_counts.index,
                                                title='Distribution of Warehouse Sizes')
                        fig_warehouses.update_layout(height=400)
                        st.plotly_chart(fig_warehouses, use_container_width=True)
                    else:
                        st.info("No warehouses defined.")
                else:
                    st.info("No warehouse sizes to display.")

            with col4:
                # Credits by Category
                if not warehouse_df.empty and 'category' in warehouse_df.columns:
                    category_credits = warehouse_df.groupby('category')['total'].sum().reset_index()
                    if not category_credits.empty and category_credits['total'].sum() > 0:
                        fig_workload = px.pie(category_credits, values='total', names='category',
                                              title='Monthly Credits by Category')
                        fig_workload.update_layout(height=400)
                        st.plotly_chart(fig_workload, use_container_width=True)
                    else:
                        st.info("No credit distribution to display.")
                else:
                    st.info("No workload categories to analyze.")
        else:
            with col3:
                st.info("No warehouse data defined.")
            with col4:
                st.info("No workload data to analyze.")

        return storage_df, compute_df, warehouse_df if warehouse_data else pd.DataFrame()

    # Create the charts
    storage_df, compute_df, warehouse_df = create_consumption_charts_from_table(consumption_data, warehouse_data, num_years)

    # Display Architecture Diagram
    display_architecture_diagram_with_png()

    # --- SAVE/UPDATE AND EXPORT UI ---
    st.markdown("---")
    st.markdown("### 💾 Save Sizing Information")
    sizing_id_to_save = st.session_state.form_data.get('sizing_id')
    if sizing_id_to_save:
        if st.button(f"✅ Update Sizing ID: {sizing_id_to_save}", key="update_in_db"):
            with st.spinner("Updating sizing data..."):
                update_in_snowflake(sizing_id_to_save)
            st.balloons()
    else:
        if st.button("💾 Save to Snowflake Database", key="save_to_db"):
            with st.spinner("Saving new sizing data..."):
                saved_id = write_to_snowflake()
                if saved_id:
                    st.balloons()
                    st.success(f"✅ Sizing information saved with new ID: **{saved_id}**")
                    st.rerun()

    # Export options
    st.markdown("### 📤 Export Options")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.file_uploader(
            "Upload PowerPoint Template (optional)",
            type="pptx",
            key="pptx_template",
            on_change=sync_widget_to_form_data,
            args=('pptx_template',)
        )
        if st.button("📊 Generate a Powerpoint Preso"):
            # (Your PowerPoint generation logic here)
            st.info("Generating PowerPoint...")
    with col2:
        if st.button("📋 Export Summary as CSV"):
            # (Your CSV export logic here)
            st.info("Generating CSV...")
    with col3:
        if st.button("📄 Generate Report"):
            st.info("Full report generation is a future feature.")

    # --- FINAL NAVIGATION ---
    st.markdown("---")
    prev_col, _, start_over_col = st.columns([1, 1, 1])
    with prev_col:
        if st.button("Previous", key="prev_6"):
            st.session_state.show_summary = False
            prev_section()
            st.rerun()
    with start_over_col:
        if st.button("Start New Sizing", key="start_over"):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()